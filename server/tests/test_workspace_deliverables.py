from io import BytesIO
from pathlib import Path

import fitz
from docx import Document as WordDocument
from pptx import Presentation

from app.documents import source_to_pdf, text_to_pdf
from app.config import Settings
from app.models import (
    AISuggestion,
    ActivityEvent,
    DocumentComment,
    DeliverableRequirement,
    DeliverableReviewFinding,
    NativeDocument,
    NativeDocumentSource,
    NativeDocumentVersion,
    Workspace,
    WorkspaceMember,
)


def assert_searchable_pdf(data: bytes, expected: str) -> None:
    assert data.startswith(b"%PDF-")
    pdf = fitz.open(stream=data, filetype="pdf")
    try:
        assert expected in "\n".join(page.get_text() for page in pdf)
    finally:
        pdf.close()


def test_plain_text_is_normalized_to_searchable_pdf() -> None:
    assert_searchable_pdf(text_to_pdf("Evidence one\nEvidence two", "Research notes"), "Evidence one")


def test_long_audit_text_is_wrapped_and_preserved_across_pdf_pages() -> None:
    requirements = "\n".join(
        f"- [x] Requirement {index}: include a detailed evidence-backed explanation with source notes and an accountable owner."
        for index in range(1, 31)
    )
    data = text_to_pdf(f"Client report\n\n# Appendix: Verification audit\n\n{requirements}\n\nRequirement coverage complete", "Verified report")
    pdf = fitz.open(stream=data, filetype="pdf")
    try:
        extracted = "\n".join(page.get_text() for page in pdf)
        assert pdf.page_count >= 2
        assert "Verification audit" in extracted
        assert "Requirement coverage complete" in extracted
    finally:
        pdf.close()


def test_docx_is_normalized_without_discarding_text() -> None:
    document = WordDocument()
    document.add_heading("Market evidence", level=1)
    document.add_paragraph("Retention improved during the pilot.")
    stream = BytesIO()
    document.save(stream)
    assert_searchable_pdf(source_to_pdf(stream.getvalue(), "docx", "brief.docx"), "Retention improved")


def test_pptx_is_normalized_without_claiming_layout_fidelity() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly finding"
    slide.placeholders[1].text = "Revenue grew in the second quarter."
    stream = BytesIO()
    presentation.save(stream)
    assert_searchable_pdf(source_to_pdf(stream.getvalue(), "pptx", "findings.pptx"), "Revenue grew")


def test_workspace_lifecycle_models_have_explicit_ownership_and_review_state() -> None:
    assert Workspace.__table__.c.owner_id.index
    assert WorkspaceMember.__table__.c.role.default.arg == "owner"
    assert NativeDocument.__table__.c.workspace_id.index
    assert NativeDocument.__table__.c.revision.default.arg == 1
    assert NativeDocumentVersion.__table__.c.native_document_id.index
    assert NativeDocumentSource.__table__.c.document_id.index
    assert DocumentComment.__table__.c.status.default.arg == "open"
    assert AISuggestion.__table__.c.status.default.arg == "pending"
    assert DeliverableRequirement.__table__.c.status.default.arg == "pending"
    assert callable(DeliverableRequirement.__table__.c.evidence.default.arg)
    assert callable(DeliverableRequirement.__table__.c.linked_sections.default.arg)
    assert DeliverableReviewFinding.__table__.c.status.default.arg == "open"
    assert DeliverableReviewFinding.__table__.c.claim_type.default.arg == "other"
    assert DeliverableReviewFinding.__table__.c.native_document_id.index
    assert ActivityEvent.__table__.c.workspace_id.index


def test_workspace_api_exposes_domain_routes() -> None:
    source = Path(__file__).parents[1].joinpath("app", "deliverables.py").read_text(encoding="utf-8")
    for route in (
        '"/workspaces"',
        '"/workspaces/{workspace_id}/search"',
        '"/workspaces/{workspace_id}/activity"',
        '"/native-documents/{native_id}/versions"',
        '"/native-documents/{native_id}/comments"',
        '"/native-documents/{native_id}/suggestions"',
        '"/native-documents/{native_id}/requirements"',
        '"/native-documents/{native_id}/requirements/extract"',
        '"/native-documents/{native_id}/review"',
        '"/native-documents/{native_id}/review-findings"',
        '"/native-documents/{native_id}/readiness"',
        '"/review-findings/{finding_id}/decision"',
        '"/native-documents/{native_id}/export"',
        '"/workspaces/{workspace_id}/demo"',
        '"/workspaces/{workspace_id}/events"',
    ):
        assert route in source


def test_verified_deliverable_migration_follows_current_head() -> None:
    migration = Path(__file__).parents[1].joinpath("alembic", "versions", "0018_verified_deliverables.py").read_text(encoding="utf-8")
    assert 'down_revision = "0017_cancellable_jobs"' in migration
    assert '"deliverable_requirements"' in migration
    assert '"deliverable_review_findings"' in migration

    repair = Path(__file__).parents[1].joinpath("alembic", "versions", "0019_normalize_generated_text.py").read_text(encoding="utf-8")
    assert 'down_revision = "0018_verified_deliverables"' in repair
    assert '"native_document_versions"' in repair
    assert "normalize_generated_text" in repair

    guided = Path(__file__).parents[1].joinpath("alembic", "versions", "0020_guided_verification.py").read_text(encoding="utf-8")
    assert 'down_revision = "0019_normalize_generated_text"' in guided
    assert '"linked_sections"' in guided
    assert '"claim_type"' in guided


def test_readiness_requires_a_draft_and_completed_verification() -> None:
    source = Path(__file__).parents[1].joinpath("app", "deliverables.py").read_text(encoding="utf-8")
    assert 'blockers.append("Write or generate the draft")' in source
    assert 'blockers.append("Run whole-deliverable verification")' in source
    assert 'status_code=409' in source
    assert '"Export is blocked until this deliverable is verified"' in source


def test_default_ai_quota_supports_a_complete_local_workflow() -> None:
    settings = Settings(_env_file=None)
    assert settings.ai_daily_request_limit >= 50
    assert settings.ai_global_daily_request_limit >= 500


def test_upload_assigns_workspace_before_document_can_autoflush() -> None:
    source = Path(__file__).parents[1].joinpath("app", "documents.py").read_text(encoding="utf-8")
    workspace_lookup = source.index("workspace = await ensure_personal_workspace(user, session)")
    document_construction = source.index("document = Document(", workspace_lookup)
    session_add = source.index("session.add(document)", document_construction)

    assert workspace_lookup < document_construction < session_add
    constructor = source[document_construction:session_add]
    assert "workspace_id=workspace.id" in constructor
