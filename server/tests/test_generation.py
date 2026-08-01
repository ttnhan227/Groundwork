from io import BytesIO

import fitz
from docx import Document
from pptx import Presentation

from app.generation import (
    CREATION_TEMPLATES,
    THEMES,
    CreateRequest,
    GeneratedContent,
    _content,
    _docx,
    _docx_dynamic,
    _pdf,
    _pdf_dynamic,
    _pptx,
    _pptx_dynamic,
    _pptx_preview_slides,
    _document_preview_pages,
    _guard_unsupported_metrics,
    _native_blocks,
)


def test_generation_builds_valid_office_and_pdf_files() -> None:
    title, sections = _content("Quarterly customer success plan", "")

    word_data = _docx(title, sections, THEMES["minimal"])
    pdf_data = _pdf(title, sections, THEMES["modern"])
    slides_data = _pptx(title, sections, THEMES["executive"])

    assert len(Document(BytesIO(word_data)).paragraphs) >= 10
    assert fitz.open(stream=pdf_data, filetype="pdf").page_count >= 1
    assert len(Presentation(BytesIO(slides_data)).slides) == 6


def test_generation_themes_cover_product_options() -> None:
    assert set(THEMES) == {"minimal", "executive", "modern", "warm"}


def test_generated_content_has_enough_sections() -> None:
    title, sections = _content("Customer retention strategy", "")
    assert title == "Customer retention strategy"
    assert len(sections) >= 5
    assert all(heading and len(body) >= 20 for heading, body in sections)


def test_dynamic_invoice_uses_table_and_distinct_layout() -> None:
    plan = GeneratedContent.model_validate({
        "title": "BrightWave Studio Invoice",
        "subtitle": "Design services for Northstar Cafe",
        "document_type": "invoice",
        "layout": "business",
        "accent_color": "#173B67",
        "metadata": [
            {"label": "Invoice number", "value": "BW-2026-041"},
            {"label": "Seller", "value": "BrightWave Studio"},
            {"label": "Client", "value": "Northstar Cafe"},
            {"label": "Due date", "value": "August 15, 2026"},
        ],
        "sections": [
            {"heading": f"Section {index}", "body": "Professional fictional invoice information for testing dynamic layouts."}
            for index in range(1, 6)
        ],
        "table": {
            "headers": ["Service", "Quantity", "Rate", "Amount"],
            "rows": [["Brand workshop", "1", "$900", "$900"], ["Menu design", "1", "$1,550", "$1,550"]],
        },
        "callout": "Subtotal $2,450 · Tax $196 · **Total $2,646**",
    })
    assert "**" not in plan.callout
    word_data = _docx_dynamic(plan, THEMES["executive"])
    pdf_data = _pdf_dynamic(plan, THEMES["executive"])
    slides_data = _pptx_dynamic(plan, THEMES["executive"])
    assert len(Document(BytesIO(word_data)).tables) >= 3
    assert fitz.open(stream=pdf_data, filetype="pdf").page_count >= 1
    assert len(Presentation(BytesIO(slides_data)).slides) == 6


def test_dynamic_presentation_constrains_long_copy_and_shapes_to_slide() -> None:
    plan = GeneratedContent.model_validate({
        "title": "A deliberately long presentation title explaining a complex transformation strategy across international markets",
        "subtitle": "A detailed subtitle for executive stakeholders that would previously collide with nearby elements when rendered in PowerPoint.",
        "document_type": "presentation",
        "layout": "modern",
        "accent_color": "#2F6BFF",
        "metadata": [],
        "sections": [
            {
                "heading": f"{index}. Operational transformation, customer outcomes, and implementation priorities",
                "body": " ".join([
                    "This slide contains deliberately verbose source material to exercise the renderer under realistic worst-case conditions.",
                    "It should remain readable, respect its allocated content region, preserve comfortable spacing, and never run beneath decorative elements.",
                    "The final presentation should communicate one focused idea instead of shrinking an entire paragraph into unreadable text.",
                ] * 2),
            }
            for index in range(1, 7)
        ],
        "table": None,
        "callout": "",
    })
    preview = _pptx_preview_slides(plan)
    assert len(preview[0]["title"]) <= 76
    assert all(len(slide["body"]) <= 330 for slide in preview[1:])

    deck = Presentation(BytesIO(_pptx_dynamic(plan, THEMES["modern"])))
    assert len(deck.slides) == 7
    for slide in deck.slides:
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0
            assert shape.left + shape.width <= deck.slide_width
            assert shape.top + shape.height <= deck.slide_height


def test_real_template_catalog_and_context_request_cover_document_workflows() -> None:
    assert len(CREATION_TEMPLATES) >= 20
    for template_id in ("annual-report", "client-proposal-document", "nda", "architecture-design", "research-paper", "quarterly-review"):
        template = CREATION_TEMPLATES[template_id]
        assert len(template["brief"]) > 80
        assert template["layout"] in {"business", "editorial", "compact", "formal", "modern"}
    request = CreateRequest.model_validate({
        "prompt": "Create the QBR",
        "output_format": "pdf",
        "template_id": "quarterly-business-review",
        "source_document_ids": ["9e501e55-352f-4c55-afcf-9412672413f4"],
        "conversation_id": "306c6a02-81d6-437f-991e-4f655f1648ce",
        "workspace_context": "Research finding: retention fell in enterprise accounts.",
        "template_answers": {"period": "Q2 2026", "decisions": "Approve retention program"},
    })
    assert request.workspace_context.startswith("Research finding")
    assert str(request.conversation_id) == "306c6a02-81d6-437f-991e-4f655f1648ce"
    assert request.template_answers["period"] == "Q2 2026"


def test_document_preview_and_native_draft_share_the_generated_structure() -> None:
    plan = GeneratedContent.model_validate({
        "title": "Q2 Business Review",
        "subtitle": "Performance, risks, and decisions",
        "document_type": "report",
        "layout": "compact",
        "accent_color": "#5B5CE2",
        "metadata": [{"label": "Period", "value": "Q2 2026"}],
        "sections": [
            {"heading": f"Section {index}", "body": "Specific evidence-backed operating content for executive review."}
            for index in range(1, 7)
        ],
        "table": {"headers": ["Metric", "Actual", "Plan"], "rows": [["Retention", "92%", "95%"]]},
        "callout": "Decision required: fund the enterprise retention program.",
    })
    pages = _document_preview_pages(plan)
    blocks = _native_blocks(plan)
    assert pages[0]["kind"] == "cover"
    assert any(page["kind"] == "table" for page in pages)
    assert len([page for page in pages if page["kind"] == "sections"]) == 3
    assert any(block["type"] == "heading" and block["text"] == "Section 1" for block in blocks)
    assert any("Retention" in block["text"] for block in blocks)


def test_unsupported_generated_metrics_are_marked_for_verification() -> None:
    plan = GeneratedContent.model_validate({
        "title": "Q2 review",
        "subtitle": "Performance update",
        "document_type": "report",
        "layout": "compact",
        "accent_color": "#5B5CE2",
        "metadata": [],
        "sections": [{"heading": "Performance", "body": "Retention improved by 22% and the 12-week rollout is complete."}],
        "table": {"headers": ["Metric", "Result"], "rows": [["Retention", "22%"]]},
        "callout": "Approve the verified $75,000 budget.",
    })
    guarded = _guard_unsupported_metrics(plan, "The approved budget is $75,000. The rollout lasted 12 weeks.")
    assert "22%" not in guarded.sections[0].body
    assert guarded.sections[0].body.count("[confirm metric]") == 1
    assert "[confirm claim]" in guarded.sections[0].body
    assert "12-week" in guarded.sections[0].body
    assert guarded.table.rows[0][1] == "[confirm metric]"
    assert "$75,000" in guarded.callout

    verified = GeneratedContent.model_validate({
        "title": "Release update", "subtitle": "", "document_type": "report", "layout": "compact",
        "accent_color": "#5B5CE2", "metadata": [],
        "sections": [{"heading": "Release", "body": "We launched the beta program."}], "table": None, "callout": "",
    })
    assert "[confirm claim]" not in _guard_unsupported_metrics(verified, "We launched the beta program in June.").sections[0].body
