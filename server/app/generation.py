import json
import re
import uuid
from io import BytesIO

import fitz
import httpx
from docx import Document as WordDocument
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from fastapi import APIRouter, Depends, HTTPException, status
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from pydantic import BaseModel, Field, field_validator
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.config import get_settings
from app.database import get_session
from app.dependencies import current_user
from app.documents import owned_document, safe_filename
from app.models import ArtifactVersion, DocumentPage, GeneratedArtifact, User
from app.schemas import ArtifactResponse
from app.storage import ObjectStorage
from app.usage import record_ai_usage

router = APIRouter(prefix="/create", tags=["Document creation"])

THEMES = {
    "minimal": {"navy": "0A183D", "accent": "EF2949", "soft": "F3F5F8"},
    "executive": {"navy": "172033", "accent": "C99A49", "soft": "F4F1EA"},
    "modern": {"navy": "14213D", "accent": "2F6BFF", "soft": "EDF3FF"},
    "warm": {"navy": "382B2A", "accent": "D86A4A", "soft": "FBF1EA"},
}

PRESENTATION_TEMPLATES = {
    "startup-pitch": {"theme": "modern", "layout": "modern", "brief": "Build a startup pitch deck: problem, solution, market, product, traction, business model, go-to-market, and next step."},
    "quarterly-review": {"theme": "executive", "layout": "compact", "brief": "Build a quarterly business review: executive summary, KPI scorecard, wins, misses, drivers, outlook, and decisions."},
    "strategy-roadmap": {"theme": "minimal", "layout": "editorial", "brief": "Build a strategy roadmap: context, strategic pillars, initiatives, phased timeline, dependencies, risks, and owners."},
    "product-launch": {"theme": "warm", "layout": "modern", "brief": "Build a product launch story: customer need, product promise, key capabilities, positioning, launch plan, channels, and success measures."},
    "data-report": {"theme": "modern", "layout": "compact", "brief": "Build a data-led report: headline findings, KPI snapshots, trends, comparisons, implications, and recommended actions."},
    "client-proposal": {"theme": "minimal", "layout": "formal", "brief": "Build a client proposal: understanding, proposed approach, workstreams, deliverables, timeline, proof, and next steps."},
}


def _plain_text(value) -> str:
    text = str(value)
    text = re.sub(r"\*\*(.*?)\*\*|__(.*?)__", lambda match: match.group(1) or match.group(2), text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    return text.strip()


class CreateRequest(BaseModel):
    prompt: str = Field(min_length=3, max_length=4000)
    output_format: str = Field(pattern="^(docx|pdf|pptx)$")
    theme: str = Field(default="minimal", pattern="^(minimal|executive|modern|warm)$")
    title: str | None = Field(default=None, max_length=160)
    source_document_id: uuid.UUID | None = None
    template_id: str | None = Field(default=None, max_length=50)


class GeneratedSection(BaseModel):
    heading: str = Field(min_length=2, max_length=90)
    body: str = Field(min_length=5, max_length=900)

    @field_validator("heading", "body", mode="before")
    @classmethod
    def stringify_section(cls, value):
        return _plain_text(value)


class GeneratedField(BaseModel):
    label: str = Field(min_length=1, max_length=50)
    value: str = Field(min_length=1, max_length=240)

    @field_validator("label", "value", mode="before")
    @classmethod
    def stringify_values(cls, value):
        return _plain_text(value)


class GeneratedTable(BaseModel):
    headers: list[str] = Field(min_length=2, max_length=5)
    rows: list[list[str]] = Field(min_length=1, max_length=20)

    @field_validator("headers", "rows", mode="before")
    @classmethod
    def stringify_table(cls, value, info):
        if info.field_name == "headers":
            return [_plain_text(item) for item in (value or [])]
        return [[_plain_text(cell) for cell in row] for row in (value or [])]

    @field_validator("rows")
    @classmethod
    def rows_match_headers(cls, rows: list[list[str]], info):
        headers = info.data.get("headers", [])
        if headers and any(len(row) != len(headers) for row in rows):
            raise ValueError("Every table row must match the header count")
        return rows


class GeneratedContent(BaseModel):
    title: str = Field(min_length=2, max_length=120)
    subtitle: str = Field(default="", max_length=180)
    document_type: str = Field(pattern="^(invoice|proposal|report|policy|resume|letter|agenda|presentation|general)$")
    layout: str = Field(pattern="^(business|editorial|compact|formal|modern)$")
    accent_color: str = Field(pattern=r"^#[0-9A-Fa-f]{6}$")
    metadata: list[GeneratedField] = Field(default_factory=list, max_length=10)
    sections: list[GeneratedSection] = Field(min_length=1, max_length=10)
    table: GeneratedTable | None = None
    callout: str = Field(default="", max_length=300)

    @field_validator("title", "subtitle", "callout", mode="before")
    @classmethod
    def clean_top_level_text(cls, value):
        return _plain_text(value)


def _clean_title(value: str) -> str:
    first = re.split(r"[.!?\n]", value.strip())[0]
    return (first[:90].strip() or "Untitled document").rstrip(":")


def _content(prompt: str, source: str) -> tuple[str, list[tuple[str, str]]]:
    title = _clean_title(prompt)
    source_sentences = [part.strip() for part in re.split(r"(?<=[.!?])\s+", source) if len(part.strip()) > 35]
    seed = source_sentences[:8]
    sections = [
        ("Executive summary", seed[0] if seed else f"This document presents a focused overview of {title.lower()}, the key considerations, and a practical path forward."),
        ("Current context", seed[1] if len(seed) > 1 else f"The opportunity is to turn {title.lower()} into a clear, audience-ready narrative supported by concrete priorities."),
        ("Key insights", seed[2] if len(seed) > 2 else "The strongest approach combines a concise message, evidence-led decisions, and a structure that makes the next action obvious."),
        ("Recommended direction", seed[3] if len(seed) > 3 else "Prioritize the highest-value outcome, define ownership, and measure progress through a small set of meaningful indicators."),
        ("Next steps", seed[4] if len(seed) > 4 else "Align stakeholders, confirm the first milestone, and begin with a focused implementation that can be reviewed and improved quickly."),
    ]
    return title, sections


async def _ai_content(payload: CreateRequest, source: str) -> GeneratedContent:
    settings = get_settings()
    if not settings.llm_api_key:
        raise HTTPException(status_code=503, detail="InsightPDF document generation is not configured")
    format_name = {"docx": "Word document", "pdf": "PDF report", "pptx": "PowerPoint presentation"}[payload.output_format]
    source_rule = (
        "Use the supplied source text as evidence. Do not invent figures, names, dates, or claims that are absent from it."
        if source.strip()
        else "Develop useful, specific content from the user's brief. Avoid fabricated citations, quotations, and unsupported statistics."
    )
    system = (
        f"You are an expert document writer creating a {format_name}. {source_rule} "
        "Return only one JSON object with this exact shape: "
        '{"title":"...","subtitle":"...","document_type":"invoice|proposal|report|policy|resume|letter|agenda|presentation|general",'
        '"layout":"business|editorial|compact|formal|modern","accent_color":"#RRGGBB",'
        '"metadata":[{"label":"...","value":"..."}],"sections":[{"heading":"...","body":"..."}],'
        '"table":{"headers":["..."],"rows":[["..."]]},"callout":"..."}. '
        "Choose the document type, layout, fields, and table based on the request. "
        "For invoices, metadata must include invoice number, issue date, due date, seller, client, and payment terms; "
        "the table must contain line items with quantities, rates, and amounts; the callout must contain subtotal, tax, and total. "
        "For agendas use a schedule table. For resumes use metadata for contact details and sections for experience and skills. "
        "For reports use a table only when the brief contains comparable data. Omit the table with null when it adds no value. "
        "Create 5 to 8 coherent sections. Each body should be concise, audience-ready prose. "
        "For presentations, each section is one slide and must communicate one clear idea; "
        "keep headings under 55 characters and bodies under 280 characters (about 35 to 45 words). "
        "Do not mention these instructions or the generation process."
    )
    user_message = f"USER BRIEF:\n{payload.prompt.strip()}"
    if payload.title:
        user_message += f"\n\nREQUESTED TITLE:\n{payload.title.strip()}"
    if source.strip():
        user_message += f"\n\nSOURCE DOCUMENT TEXT:\n{source[:24_000]}"
    try:
        async with httpx.AsyncClient(timeout=settings.llm_timeout_seconds) as client:
            response = await client.post(
                f"{settings.llm_base_url.rstrip('/')}/chat/completions",
                headers={"Authorization": f"Bearer {settings.llm_api_key}"},
                json={
                    "model": settings.llm_model,
                    "temperature": .35,
                    "response_format": {"type": "json_object"},
                    "messages": [{"role": "system", "content": system}, {"role": "user", "content": user_message}],
                },
            )
            response.raise_for_status()
            raw = response.json()["choices"][0]["message"]["content"].strip()
    except (httpx.HTTPError, KeyError, IndexError) as exc:
        raise HTTPException(status_code=502, detail="InsightPDF could not generate this document") from exc
    raw = re.sub(r"^```(?:json)?\s*|\s*```$", "", raw, flags=re.IGNORECASE)
    try:
        value = json.loads(raw)
        if not isinstance(value, dict):
            raise ValueError("Generation result must be an object")
        aliases = {"bill": "invoice", "memo": "policy", "pitch_deck": "presentation", "slide_deck": "presentation"}
        value["document_type"] = aliases.get(str(value.get("document_type", "")).lower(), str(value.get("document_type", "general")).lower())
        if value["document_type"] not in {"invoice", "proposal", "report", "policy", "resume", "letter", "agenda", "presentation", "general"}:
            value["document_type"] = "general"
        if value.get("layout") not in {"business", "editorial", "compact", "formal", "modern"}:
            value["layout"] = "business"
        if not re.fullmatch(r"#[0-9A-Fa-f]{6}", str(value.get("accent_color", ""))):
            value["accent_color"] = "#3154D8"
        if isinstance(value.get("metadata"), dict):
            value["metadata"] = [{"label": key.replace("_", " ").title(), "value": item} for key, item in value["metadata"].items()]
        if isinstance(value.get("callout"), dict):
            value["callout"] = " · ".join(f"{key.replace('_', ' ').title()}: {item}" for key, item in value["callout"].items())
        table = value.get("table")
        if isinstance(table, dict) and isinstance(table.get("headers"), list) and isinstance(table.get("rows"), list):
            width = len(table["headers"])
            table["rows"] = [(list(row) + [""] * width)[:width] for row in table["rows"] if isinstance(row, list)]
            if not table["rows"]:
                value["table"] = None
        result = GeneratedContent.model_validate(value)
    except (json.JSONDecodeError, ValueError) as exc:
        raise HTTPException(status_code=502, detail="InsightPDF returned an invalid document structure") from exc
    return result


def _docx(title: str, sections: list[tuple[str, str]], theme: dict[str, str]) -> bytes:
    document = WordDocument()
    section = document.sections[0]
    section.top_margin = Inches(.72)
    section.bottom_margin = Inches(.72)
    section.left_margin = Inches(.82)
    section.right_margin = Inches(.82)
    heading = document.add_heading(title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    heading.runs[0].font.color.rgb = RGBColor.from_string(theme["navy"])
    heading.runs[0].font.size = Pt(28)
    sub = document.add_paragraph("Created with InsightPDF AI")
    sub.runs[0].font.color.rgb = RGBColor.from_string(theme["accent"])
    sub.runs[0].font.size = Pt(10)
    for index, (name, body) in enumerate(sections, 1):
        heading = document.add_heading(f"{index:02d}  {name}", level=1)
        heading.runs[0].font.color.rgb = RGBColor.from_string(theme["navy"])
        paragraph = document.add_paragraph(body)
        paragraph.paragraph_format.space_after = Pt(12)
        paragraph.paragraph_format.line_spacing = 1.18
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pdf(title: str, sections: list[tuple[str, str]], theme: dict[str, str]) -> bytes:
    pdf = fitz.open()
    page = pdf.new_page()
    accent = tuple(int(theme["accent"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    navy = tuple(int(theme["navy"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    page.draw_rect(fitz.Rect(0, 0, 14, page.rect.height), color=accent, fill=accent)
    page.insert_textbox(fitz.Rect(48, 58, 545, 150), title, fontsize=26, fontname="hebo", color=navy)
    page.insert_text((48, 158), "Created with InsightPDF AI", fontsize=9, color=accent)
    y = 205
    for index, (name, body) in enumerate(sections, 1):
        page.insert_text((48, y), f"{index:02d}  {name.upper()}", fontsize=11, fontname="hebo", color=navy)
        y += 18
        used = page.insert_textbox(fitz.Rect(48, y, 545, y + 72), body, fontsize=10, lineheight=1.25, color=(.26, .3, .38))
        y += max(60, 78 - min(used, 0))
        if y > 730 and index < len(sections):
            page = pdf.new_page()
            y = 60
    return pdf.tobytes(garbage=4, deflate=True)


def _pptx(title: str, sections: list[tuple[str, str]], theme: dict[str, str]) -> bytes:
    deck = Presentation()
    deck.slide_width = PptxInches(13.333)
    deck.slide_height = PptxInches(7.5)
    colors = {key: PptxRGBColor.from_string(value) for key, value in theme.items()}

    def add_text(slide, text, left, top, width, height, size, color, bold=False):
        box = slide.shapes.add_textbox(PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = "Aptos"
        paragraph.font.size = PptxPt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        return box

    title_slide = deck.slides.add_slide(deck.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = colors["navy"]
    add_text(title_slide, "INSIGHTPDF PRESENTATION", .85, .75, 5, .35, 13, colors["accent"], True)
    add_text(title_slide, title, .85, 2.0, 9.4, 2.0, 34, PptxRGBColor(255, 255, 255), True)
    add_text(title_slide, "A clear, AI-created narrative", .85, 5.75, 6, .4, 18, PptxRGBColor(190, 200, 218))
    for index, (name, body) in enumerate(sections, 1):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptxRGBColor(255, 255, 255)
        accent = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(.16), deck.slide_height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = colors["accent"]
        accent.line.fill.background()
        add_text(slide, f"{index:02d} · {name.upper()}", .82, .72, 5, .35, 13, colors["accent"], True)
        add_text(slide, name, .82, 1.45, 10.7, .85, 30, colors["navy"], True)
        add_text(slide, body, .82, 2.65, 8.9, 2.2, 20, PptxRGBColor(62, 72, 91))
        add_text(slide, "InsightPDF", .82, 6.7, 2, .3, 10, PptxRGBColor(130, 140, 158), True)
        add_text(slide, f"{index + 1:02d}", 11.8, 6.7, .5, .3, 10, PptxRGBColor(130, 140, 158))
    output = BytesIO()
    deck.save(output)
    return output.getvalue()


def _shade_cell(cell, color: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    fill = OxmlElement("w:shd")
    fill.set(qn("w:fill"), color)
    properties.append(fill)


def _docx_dynamic(plan: GeneratedContent, theme: dict[str, str]) -> bytes:
    document = WordDocument()
    section = document.sections[0]
    section.top_margin = Inches(.68)
    section.bottom_margin = Inches(.68)
    section.left_margin = Inches(.78)
    section.right_margin = Inches(.78)
    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5 if plan.layout == "compact" else 11)
    for style_name, size in (("Heading 1", 16), ("Heading 2", 13)):
        styles[style_name].font.name = "Aptos Display"
        styles[style_name].font.size = Pt(size)
        styles[style_name].font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))

    if plan.document_type == "invoice":
        brand = document.add_paragraph()
        brand.paragraph_format.space_after = Pt(3)
        run = brand.add_run(next((item.value for item in plan.metadata if item.label.lower() in {"seller", "from"}), "INVOICE"))
        run.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor.from_string(theme["navy"])
        title = document.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        title.paragraph_format.space_after = Pt(16)
        title_run = title.add_run("INVOICE")
        title_run.bold = True
        title_run.font.size = Pt(30)
        title_run.font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))
        details = document.add_table(rows=0, cols=2)
        details.alignment = WD_TABLE_ALIGNMENT.CENTER
        details.autofit = False
        for item in plan.metadata:
            cells = details.add_row().cells
            cells[0].width, cells[1].width = Inches(1.55), Inches(4.95)
            cells[0].text, cells[1].text = item.label.upper(), item.value
            cells[0].paragraphs[0].runs[0].font.bold = True
            cells[0].paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))
        document.add_paragraph()
    else:
        kicker = document.add_paragraph(plan.document_type.upper())
        kicker.paragraph_format.space_after = Pt(6)
        kicker.runs[0].bold = True
        kicker.runs[0].font.size = Pt(9)
        kicker.runs[0].font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))
        heading = document.add_paragraph()
        heading.paragraph_format.space_after = Pt(5)
        heading_run = heading.add_run(plan.title)
        heading_run.bold = True
        heading_run.font.size = Pt(28 if plan.layout == "editorial" else 24)
        heading_run.font.color.rgb = RGBColor.from_string(theme["navy"])
        if plan.subtitle:
            subtitle = document.add_paragraph(plan.subtitle)
            subtitle.paragraph_format.space_after = Pt(15)
            subtitle.runs[0].italic = True
            subtitle.runs[0].font.color.rgb = RGBColor(100, 110, 126)
        if plan.metadata:
            metadata = document.add_table(rows=0, cols=2)
            metadata.alignment = WD_TABLE_ALIGNMENT.CENTER
            for index in range(0, len(plan.metadata), 2):
                cells = metadata.add_row().cells
                for offset, cell in enumerate(cells):
                    if index + offset < len(plan.metadata):
                        item = plan.metadata[index + offset]
                        cell.text = f"{item.label.upper()}\n{item.value}"
                        cell.paragraphs[0].runs[0].font.bold = True
                        _shade_cell(cell, theme["soft"])
            document.add_paragraph()

    if plan.table:
        table = document.add_table(rows=1, cols=len(plan.table.headers))
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.style = "Table Grid"
        table.autofit = True
        for index, header in enumerate(plan.table.headers):
            cell = table.rows[0].cells[index]
            cell.text = header
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            _shade_cell(cell, plan.accent_color.removeprefix("#"))
            for run in cell.paragraphs[0].runs:
                run.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
        for row in plan.table.rows:
            cells = table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = value
                cells[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                if len(table.rows) % 2 == 1:
                    _shade_cell(cells[index], theme["soft"])
        document.add_paragraph()

    if plan.callout:
        callout = document.add_table(rows=1, cols=1)
        callout.alignment = WD_TABLE_ALIGNMENT.RIGHT if plan.document_type == "invoice" else WD_TABLE_ALIGNMENT.CENTER
        callout.cell(0, 0).text = plan.callout
        _shade_cell(callout.cell(0, 0), theme["soft"])
        callout.cell(0, 0).paragraphs[0].runs[0].font.bold = True
        callout.cell(0, 0).paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(theme["navy"])
        document.add_paragraph()

    visible_sections = plan.sections[:2] if plan.document_type == "invoice" else plan.sections
    for index, item in enumerate(visible_sections, 1):
        heading = document.add_heading(item.heading, level=1)
        heading.paragraph_format.keep_with_next = True
        body = document.add_paragraph(item.body)
        body.paragraph_format.space_after = Pt(9)
        body.paragraph_format.line_spacing = 1.15 if plan.layout != "editorial" else 1.28
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run(f"Created with InsightPDF AI · {plan.document_type.title()} · {plan.layout.title()} layout").font.size = Pt(8)
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pdf_dynamic(plan: GeneratedContent, theme: dict[str, str]) -> bytes:
    document = fitz.open()
    accent = tuple(int(plan.accent_color[index:index + 2], 16) / 255 for index in (1, 3, 5))
    navy = tuple(int(theme["navy"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    soft = tuple(int(theme["soft"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    page = document.new_page()
    page.draw_rect(fitz.Rect(0, 0, 12, page.rect.height), color=accent, fill=accent)
    if plan.document_type == "invoice":
        page.insert_text((48, 66), "INVOICE", fontsize=30, fontname="hebo", color=accent)
        page.insert_text((48, 92), plan.title, fontsize=11, color=navy)
    else:
        page.insert_text((48, 54), plan.document_type.upper(), fontsize=9, fontname="hebo", color=accent)
        page.insert_textbox(fitz.Rect(48, 72, 545, 145), plan.title, fontsize=25, fontname="hebo", color=navy)
        if plan.subtitle:
            page.insert_textbox(fitz.Rect(48, 140, 545, 180), plan.subtitle, fontsize=10, color=(.38, .42, .5))
    y = 130 if plan.document_type == "invoice" else 190
    if plan.metadata:
        for index, item in enumerate(plan.metadata):
            column = index % 2
            row = index // 2
            x = 48 + column * 250
            py = y + row * 35
            page.insert_text((x, py), item.label.upper(), fontsize=7, fontname="hebo", color=accent)
            page.insert_textbox(fitz.Rect(x, py + 5, x + 225, py + 30), item.value, fontsize=9, color=navy)
        y += ((len(plan.metadata) + 1) // 2) * 35 + 15
    if plan.table:
        column_width = 497 / len(plan.table.headers)
        page.draw_rect(fitz.Rect(48, y, 545, y + 25), color=accent, fill=accent)
        for index, header in enumerate(plan.table.headers):
            page.insert_textbox(fitz.Rect(52 + index * column_width, y + 7, 48 + (index + 1) * column_width - 4, y + 23), header, fontsize=8, fontname="hebo", color=(1, 1, 1))
        y += 25
        for row_index, row in enumerate(plan.table.rows):
            if y > 720:
                page = document.new_page()
                y = 55
            if row_index % 2 == 0:
                page.draw_rect(fitz.Rect(48, y, 545, y + 27), color=soft, fill=soft)
            for index, value in enumerate(row):
                page.insert_textbox(fitz.Rect(52 + index * column_width, y + 7, 48 + (index + 1) * column_width - 4, y + 26), value, fontsize=8, color=navy)
            y += 27
        y += 12
    if plan.callout:
        callout_text = plan.callout.replace(" · ", "\n")
        callout_height = 66 if plan.document_type == "invoice" else 48
        page.draw_rect(fitz.Rect(315, y, 545, y + callout_height), color=soft, fill=soft)
        page.insert_textbox(fitz.Rect(328, y + 9, 533, y + callout_height - 7), callout_text, fontsize=9, fontname="hebo", lineheight=1.15, align=2, color=navy)
        y += callout_height + 18
    if plan.document_type == "invoice":
        for index, item in enumerate(plan.sections[:2]):
            x = 48 + index * 250
            page.insert_text((x, y), item.heading.upper(), fontsize=9, fontname="hebo", color=accent)
            page.insert_textbox(fitz.Rect(x, y + 12, x + 225, min(page.rect.height - 42, y + 105)), item.body, fontsize=8.5, lineheight=1.18, color=(.25, .29, .36))
        return document.tobytes(garbage=4, deflate=True)
    for item in plan.sections:
        if y > 700:
            page = document.new_page()
            y = 58
        page.insert_text((48, y), item.heading.upper(), fontsize=10, fontname="hebo", color=accent)
        y += 15
        page.insert_textbox(fitz.Rect(48, y, 545, y + 74), item.body, fontsize=9.5, lineheight=1.25, color=(.25, .29, .36))
        y += 77
    return document.tobytes(garbage=4, deflate=True)


def _shorten_slide_text(value: str, limit: int) -> str:
    value = re.sub(r"\s+", " ", value).strip()
    if len(value) <= limit:
        return value
    candidate = value[:limit - 1].rsplit(" ", 1)[0].rstrip(" ,;:-")
    return f"{candidate or value[:limit - 1].rstrip()}…"


def _pptx_preview_slides(plan: GeneratedContent) -> list[dict[str, str]]:
    slides = [{
        "eyebrow": _shorten_slide_text(plan.document_type.upper(), 34),
        "title": _shorten_slide_text(plan.title, 76),
        "body": _shorten_slide_text(plan.subtitle or "Created with InsightPDF AI", 130),
        "variant": "title",
    }]
    limits = {1: (64, 330), 2: (48, 250), 3: (68, 310)}
    for index, section in enumerate(plan.sections, 1):
        variant = ((index - 1) % 3) + 1
        heading_limit, body_limit = limits[variant]
        slides.append({
            "eyebrow": f"{index:02d}",
            "title": _shorten_slide_text(section.heading, heading_limit),
            "body": _shorten_slide_text(section.body, body_limit),
            "variant": f"section-{variant}",
        })
    return slides


def _pptx_dynamic(plan: GeneratedContent, theme: dict[str, str], hero_image: bytes | None = None) -> bytes:
    slide_data = _pptx_preview_slides(plan)
    sections = [(item["title"], item["body"]) for item in slide_data[1:]]
    deck = Presentation()
    deck.slide_width = PptxInches(13.333)
    deck.slide_height = PptxInches(7.5)
    navy = PptxRGBColor.from_string(theme["navy"])
    accent = PptxRGBColor.from_string(plan.accent_color.removeprefix("#"))
    soft = PptxRGBColor.from_string(theme["soft"])

    def text(slide, value, left, top, width, height, size, color, bold=False,
             align=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP):
        shape = slide.shapes.add_textbox(PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height))
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = vertical
        paragraph = frame.paragraphs[0]
        paragraph.text = value
        paragraph.alignment = align
        paragraph.space_before = paragraph.space_after = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = PptxPt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        return shape

    title_slide = deck.slides.add_slide(deck.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = navy
    title_data = slide_data[0]
    if hero_image:
        try:
            title_slide.shapes.add_picture(BytesIO(hero_image), PptxInches(8.7), PptxInches(0), width=PptxInches(4.633), height=PptxInches(7.5))
        except Exception:
            hero_image = None
    text(title_slide, title_data["eyebrow"], .85, .72, 5, .3, 14, accent, True)
    text(title_slide, title_data["title"], .85, 1.65, 7.35 if hero_image else 11.25, 2.25, 50, PptxRGBColor(255, 255, 255), True,
         vertical=MSO_ANCHOR.MIDDLE)
    text(title_slide, title_data["body"], .85, 5.55, 7.15 if hero_image else 9.5, .85, 19, PptxRGBColor(190, 200, 218))
    for index, (heading, body) in enumerate(sections, 1):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptxRGBColor(255, 255, 255) if index % 3 else navy
        ink = navy if index % 3 else PptxRGBColor(255, 255, 255)
        muted = PptxRGBColor(62, 72, 91) if index % 3 else PptxRGBColor(195, 204, 219)
        if index % 3 == 1:
            text(slide, f"{index:02d}", .8, .65, 1, .35, 18, accent, True)
            text(slide, heading, .8, 1.35, 9.25, 1.15, 35, ink, True, vertical=MSO_ANCHOR.MIDDLE)
            text(slide, body, .8, 2.95, 8.95, 2.55, 19, muted)
            marker = slide.shapes.add_shape(1, PptxInches(10.75), PptxInches(1.35), PptxInches(1.65), PptxInches(4.85))
            marker.fill.solid(); marker.fill.fore_color.rgb = accent; marker.line.fill.background()
        elif index % 3 == 2:
            panel = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(4.9), deck.slide_height)
            panel.fill.solid(); panel.fill.fore_color.rgb = soft; panel.line.fill.background()
            text(slide, f"{index:02d}", 1.05, 2.35, 2.7, 1.4, 60, accent, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            text(slide, heading, 5.75, 1.2, 6.55, 1.45, 35, ink, True, vertical=MSO_ANCHOR.MIDDLE)
            text(slide, body, 5.75, 3.0, 6.35, 2.4, 18, muted)
        else:
            text(slide, f"{index:02d}  /  INSIGHT", .85, .78, 4, .3, 14, accent, True)
            text(slide, heading, .85, 1.75, 11.45, 1.35, 37, ink, True, vertical=MSO_ANCHOR.MIDDLE)
            text(slide, body, .85, 3.5, 10.9, 2.25, 19, muted)
        text(slide, "InsightPDF", .8, 6.8, 2, .25, 10, muted, True)
    output = BytesIO()
    deck.save(output)
    return output.getvalue()


@router.get("/themes")
async def themes() -> list[dict[str, str]]:
    return [
        {"id": "minimal", "name": "Minimal", "description": "Crisp navy and red"},
        {"id": "executive", "name": "Executive", "description": "Charcoal and gold"},
        {"id": "modern", "name": "Modern", "description": "Bold blue and white"},
        {"id": "warm", "name": "Warm", "description": "Editorial earth tones"},
    ]


@router.post("", response_model=ArtifactResponse, status_code=status.HTTP_201_CREATED)
async def create_file(
    payload: CreateRequest,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> GeneratedArtifact:
    source_text = ""
    if payload.source_document_id:
        await owned_document(payload.source_document_id, user, session)
        pages = list(await session.scalars(
            select(DocumentPage).where(DocumentPage.document_id == payload.source_document_id).order_by(DocumentPage.page_number)
        ))
        source_text = " ".join(page.text for page in pages)[:24_000]
    template = PRESENTATION_TEMPLATES.get(payload.template_id or "")
    content_request = payload
    if payload.output_format == "pptx" and template:
        content_request = payload.model_copy(update={"prompt": f"{template['brief']}\n\nTOPIC AND REQUIREMENTS:\n{payload.prompt}"})
    plan = await _ai_content(content_request, source_text)
    await record_ai_usage(user, f"create_{payload.output_format}", session, cached=False)
    selected_theme = template["theme"] if payload.output_format == "pptx" and template else payload.theme
    if payload.output_format == "pptx" and template:
        plan.layout = template["layout"]
    theme = THEMES[selected_theme]
    data = _pptx_dynamic(plan, theme) if payload.output_format == "pptx" else {"docx": _docx_dynamic, "pdf": _pdf_dynamic}[payload.output_format](plan, theme)
    filename = safe_filename(re.sub(r"[^a-zA-Z0-9 -]", "", plan.title).strip().replace(" ", "-").lower()[:70] or "insight-document")
    filename = f"{filename}.{payload.output_format}"
    content_types = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pdf": "application/pdf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    identifier = uuid.uuid4()
    key = f"{user.id}/generated/{identifier}/{filename}"
    try:
        ObjectStorage().upload(key, data, content_types[payload.output_format])
    except Exception as exc:
        raise HTTPException(status_code=503, detail="Generated-file storage is temporarily unavailable") from exc
    artifact = GeneratedArtifact(
        id=identifier,
        owner_id=user.id,
        operation=f"ai_create_{payload.output_format}",
        filename=filename,
        object_key=key,
        content_type=content_types[payload.output_format],
        size_bytes=len(data),
        parameters={
            "prompt": payload.prompt,
            "theme": selected_theme,
            "template_id": payload.template_id,
            "source_document_id": str(payload.source_document_id) if payload.source_document_id else None,
            "provider": "mistral",
            "model": get_settings().llm_model,
            "document_type": plan.document_type,
            "layout": plan.layout,
            "accent_color": plan.accent_color,
            "preview_palette": {
                "navy": f"#{theme['navy']}",
                "accent": plan.accent_color,
                "soft": f"#{theme['soft']}",
            },
            "preview_slides": _pptx_preview_slides(plan) if payload.output_format == "pptx" else [],
        },
    )
    session.add(artifact)
    session.add(ArtifactVersion(
        artifact_id=identifier,
        version_number=1,
        object_key=key,
        content_type=content_types[payload.output_format],
        size_bytes=len(data),
        change_prompt=payload.prompt,
        metadata_json={"operation": artifact.operation, "theme": selected_theme},
    ))
    await session.commit()
    await session.refresh(artifact)
    return artifact
