import re
import uuid
from io import BytesIO

import fitz
from docx import Document as WordDocument
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from fastapi import APIRouter, Depends, HTTPException, status
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from pydantic import BaseModel, Field, field_validator
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.ai_orchestration import AIProviderError, ai_orchestrator
from app.config import get_settings
from app.database import get_session
from app.dependencies import current_user
from app.documents import owned_document, safe_filename
from app.models import (
    ArtifactVersion,
    DocumentPage,
    GeneratedArtifact,
    NativeDocument,
    NativeDocumentSource,
    NativeDocumentVersion,
    User,
)
from app.schemas import ArtifactResponse, ProcessingJobResponse
from app.storage import ObjectStorage
from app.usage import record_ai_usage

router = APIRouter(prefix="/create", tags=["Document creation"])

THEMES = {
    "minimal": {"navy": "0A183D", "accent": "EF2949", "soft": "F3F5F8"},
    "executive": {"navy": "172033", "accent": "C99A49", "soft": "F4F1EA"},
    "modern": {"navy": "14213D", "accent": "2F6BFF", "soft": "EDF3FF"},
    "warm": {"navy": "382B2A", "accent": "D86A4A", "soft": "FBF1EA"},
}

PRESENTATION_TEMPLATES = {
    "startup-pitch": {"theme": "modern", "layout": "modern", "brief": "Build a startup pitch deck: problem, solution, market, product, traction, business model, go-to-market, and next step."},
    "quarterly-review": {"theme": "executive", "layout": "compact", "brief": "Build a quarterly business review: executive summary, KPI scorecard, wins, misses, drivers, outlook, and decisions."},
    "strategy-roadmap": {"theme": "minimal", "layout": "editorial", "brief": "Build a strategy roadmap: context, strategic pillars, initiatives, phased timeline, dependencies, risks, and owners."},
    "product-launch": {"theme": "warm", "layout": "modern", "brief": "Build a product launch story: customer need, product promise, key capabilities, positioning, launch plan, channels, and success measures."},
    "data-report": {"theme": "modern", "layout": "compact", "brief": "Build a data-led report: headline findings, KPI snapshots, trends, comparisons, implications, and recommended actions."},
    "client-proposal": {"theme": "minimal", "layout": "formal", "brief": "Build a client proposal: understanding, proposed approach, workstreams, deliverables, timeline, proof, and next steps."},
}

DOCUMENT_TEMPLATES = {
    "verified-client-report": {"theme": "modern", "layout": "editorial", "type": "report", "brief": "Create a client-ready report driven by the engagement brief: executive summary, explicit acceptance criteria, evidence-backed situation analysis, key findings, recommendations, implementation plan, risks and assumptions, source notes, and a requirement-coverage appendix. Mark missing inputs clearly and never invent evidence."},
    "annual-report": {"theme": "executive", "layout": "editorial", "type": "report", "brief": "Create an annual report with a letter from leadership, year in review, performance scorecard, operating highlights, risks, outlook, and priorities for the next year."},
    "quarterly-business-review": {"theme": "modern", "layout": "compact", "type": "report", "brief": "Create a quarterly business review with an executive readout, KPI scorecard, wins, misses, customer signals, operating risks, decisions required, and next-quarter commitments."},
    "incident-report": {"theme": "minimal", "layout": "business", "type": "report", "brief": "Create a production incident report with severity and ownership metadata, impact, detection, precise timeline, root cause, contributing factors, response assessment, corrective actions, and prevention measures."},
    "audit-report": {"theme": "executive", "layout": "formal", "type": "report", "brief": "Create a formal audit report with scope, methodology, control assessment, evidence-backed findings, severity ratings, management responses, remediation owners, and due dates."},
    "research-report": {"theme": "minimal", "layout": "editorial", "type": "report", "brief": "Create a rigorous research report with research question, methodology, evidence synthesis, findings, limitations, implications, recommendations, and source notes."},
    "client-proposal-document": {"theme": "warm", "layout": "modern", "type": "proposal", "brief": "Create a persuasive client proposal with an understanding of the client's situation, objectives, proposed approach, workstreams, deliverables, timeline, team, commercial assumptions, success measures, and next steps."},
    "product-requirements-document": {"theme": "modern", "layout": "business", "type": "general", "brief": "Create a product requirements document with problem statement, users and jobs, goals and non-goals, scope, user stories, functional requirements, UX principles, analytics, dependencies, risks, rollout, and acceptance criteria."},
    "business-plan": {"theme": "executive", "layout": "editorial", "type": "proposal", "brief": "Create an investor-ready business plan with company thesis, market, customer, product, business model, competition, go-to-market, operations, milestones, financial assumptions, risks, and funding use."},
    "status-report": {"theme": "minimal", "layout": "compact", "type": "report", "brief": "Create a concise status report with overall health, progress against plan, completed work, milestones, metrics, blockers, risks, decisions, owners, and next reporting-period priorities."},
    "meeting-minutes": {"theme": "minimal", "layout": "compact", "type": "agenda", "brief": "Create decision-oriented meeting minutes with attendees, objective, discussion summary by topic, decisions and rationale, action items with owners and due dates, open questions, and next meeting."},
    "offer-letter": {"theme": "warm", "layout": "formal", "type": "letter", "brief": "Create a professional offer letter with role, reporting line, location, start date, compensation, benefits, conditions, confidentiality reference, acceptance deadline, and signature blocks. Clearly mark any missing legal details for review."},
    "employee-handbook": {"theme": "modern", "layout": "business", "type": "policy", "brief": "Create an employee handbook with company principles, employment basics, workplace conduct, compensation and benefits, time off, security, acceptable use, reporting channels, acknowledgements, and policy ownership."},
    "performance-review": {"theme": "warm", "layout": "business", "type": "report", "brief": "Create a balanced performance review with role expectations, outcomes, competency evidence, strengths, development areas, feedback themes, rating rationale, goals, and manager support actions."},
    "nda": {"theme": "executive", "layout": "formal", "type": "policy", "brief": "Create a mutual non-disclosure agreement draft with parties, purpose, definition of confidential information, exclusions, permitted use, safeguards, compelled disclosure, term, return or destruction, remedies, and governing law. Flag jurisdiction-specific clauses for legal review."},
    "service-agreement": {"theme": "executive", "layout": "formal", "type": "policy", "brief": "Create a professional services agreement draft with parties, services, deliverables, fees, change control, responsibilities, acceptance, IP, confidentiality, warranties, liability, termination, disputes, and signatures. Flag legal-review items."},
    "privacy-policy": {"theme": "minimal", "layout": "formal", "type": "policy", "brief": "Create a plain-language privacy policy with scope, data collected, purposes, legal bases where relevant, sharing, retention, security, international transfers, user rights, children, changes, and contact details. Flag jurisdiction-specific requirements."},
    "software-requirements-specification": {"theme": "modern", "layout": "business", "type": "general", "brief": "Create a software requirements specification with system purpose, actors, assumptions, functional requirements with identifiers, interfaces, data requirements, non-functional requirements, constraints, traceability, and acceptance criteria."},
    "architecture-design": {"theme": "modern", "layout": "editorial", "type": "general", "brief": "Create an architecture design document with context, requirements and quality attributes, system boundaries, component responsibilities, data flows, interfaces, security, deployment, observability, failure modes, trade-offs, decisions, and rollout."},
    "api-documentation": {"theme": "minimal", "layout": "compact", "type": "general", "brief": "Create usable API documentation with overview, authentication, environments, conventions, endpoint reference, request and response examples, errors, pagination, rate limits, idempotency, security, and versioning."},
    "test-plan": {"theme": "minimal", "layout": "business", "type": "general", "brief": "Create a release test plan with objectives, scope and exclusions, quality risks, environments, test data, test types, scenarios, entry and exit criteria, defect process, responsibilities, schedule, and reporting."},
    "runbook": {"theme": "modern", "layout": "compact", "type": "general", "brief": "Create an operational runbook with service overview, ownership, dependencies, access prerequisites, health checks, alerts, standard procedures, incident triage, recovery, rollback, escalation, verification, and maintenance."},
    "research-paper": {"theme": "minimal", "layout": "editorial", "type": "report", "brief": "Create an academic research paper structure with abstract, research question, related context, methodology, results, discussion, limitations, conclusion, and references placeholders. Never invent citations."},
}

CREATION_TEMPLATES = {**DOCUMENT_TEMPLATES, **PRESENTATION_TEMPLATES}


def _plain_text(value) -> str:
    text = str(value)
    text = re.sub(r"\*\*(.*?)\*\*|__(.*?)__", lambda match: match.group(1) or match.group(2), text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    return text.strip()


class CreateRequest(BaseModel):
    prompt: str = Field(min_length=3, max_length=4000)
    output_format: str = Field(pattern="^(docx|pdf|pptx)$")
    theme: str = Field(default="minimal", pattern="^(minimal|executive|modern|warm)$")
    title: str | None = Field(default=None, max_length=160)
    source_document_id: uuid.UUID | None = None
    source_document_ids: list[uuid.UUID] = Field(default_factory=list, max_length=12)
    conversation_id: uuid.UUID | None = None
    template_id: str | None = Field(default=None, max_length=50)
    workspace_context: str = Field(default="", max_length=12000)
    template_answers: dict[str, str] = Field(default_factory=dict)

    @field_validator("template_answers", mode="before")
    @classmethod
    def clean_template_answers(cls, value):
        if not isinstance(value, dict):
            return {}
        return {
            _plain_text(key)[:80]: _plain_text(item)[:1000]
            for key, item in list(value.items())[:16]
            if _plain_text(key) and _plain_text(item)
        }


class GeneratedSection(BaseModel):
    heading: str = Field(min_length=2, max_length=90)
    body: str = Field(min_length=5, max_length=900)

    @field_validator("heading", "body", mode="before")
    @classmethod
    def stringify_section(cls, value):
        return _plain_text(value)


class GeneratedField(BaseModel):
    label: str = Field(min_length=1, max_length=50)
    value: str = Field(min_length=1, max_length=240)

    @field_validator("label", "value", mode="before")
    @classmethod
    def stringify_values(cls, value):
        return _plain_text(value)


class GeneratedTable(BaseModel):
    headers: list[str] = Field(min_length=2, max_length=5)
    rows: list[list[str]] = Field(min_length=1, max_length=20)

    @field_validator("headers", "rows", mode="before")
    @classmethod
    def stringify_table(cls, value, info):
        if info.field_name == "headers":
            return [_plain_text(item) for item in (value or [])]
        return [[_plain_text(cell) for cell in row] for row in (value or [])]

    @field_validator("rows")
    @classmethod
    def rows_match_headers(cls, rows: list[list[str]], info):
        headers = info.data.get("headers", [])
        if headers and any(len(row) != len(headers) for row in rows):
            raise ValueError("Every table row must match the header count")
        return rows


class GeneratedContent(BaseModel):
    title: str = Field(min_length=2, max_length=120)
    subtitle: str = Field(default="", max_length=180)
    document_type: str = Field(pattern="^(invoice|proposal|report|policy|resume|letter|agenda|presentation|general)$")
    layout: str = Field(pattern="^(business|editorial|compact|formal|modern)$")
    accent_color: str = Field(pattern=r"^#[0-9A-Fa-f]{6}$")
    metadata: list[GeneratedField] = Field(default_factory=list, max_length=10)
    sections: list[GeneratedSection] = Field(min_length=1, max_length=10)
    table: GeneratedTable | None = None
    callout: str = Field(default="", max_length=300)

    @field_validator("title", "subtitle", "callout", mode="before")
    @classmethod
    def clean_top_level_text(cls, value):
        return _plain_text(value)


def _clean_title(value: str) -> str:
    first = re.split(r"[.!?\n]", value.strip())[0]
    return (first[:90].strip() or "Untitled document").rstrip(":")


_METRIC_PATTERN = re.compile(
    r"(?<!\w)(?:[$€£]\s?\d[\d,.]*|\d+(?:\.\d+)?\s?(?:%|percent|million|billion|thousand|users|customers|hours|days|weeks|months|years))(?!\w)",
    re.IGNORECASE,
)
_COMPLETED_CLAIM_PATTERN = re.compile(
    r"\b(launched|introduced|implemented|deployed|overhauled|achieved|reported|grew|increased|reduced|improved|completed|migrated|expanded|delivered|adopted|reached)\b",
    re.IGNORECASE,
)


def _guard_unsupported_metrics(plan: GeneratedContent, evidence: str) -> GeneratedContent:
    """Replace model-invented quantitative claims with explicit verification markers."""
    normalized_evidence = re.sub(r"\s+", " ", evidence).lower()

    def guard(value: str) -> str:
        def replacement(match: re.Match) -> str:
            metric = re.sub(r"\s+", " ", match.group(0)).lower()
            return match.group(0) if metric in normalized_evidence else "[confirm metric]"
        guarded = _METRIC_PATTERN.sub(replacement, value)
        sentences = re.split(r"(?<=[.!?])\s+", guarded)
        for index, sentence in enumerate(sentences):
            claim = _COMPLETED_CLAIM_PATTERN.search(sentence)
            if claim and claim.group(1).lower() not in normalized_evidence and "[confirm claim]" not in sentence.lower():
                sentences[index] = f"{sentence.rstrip()} [confirm claim]"
        return " ".join(sentences)

    plan.title = guard(plan.title)
    plan.subtitle = guard(plan.subtitle)
    plan.callout = guard(plan.callout)
    for item in plan.metadata:
        item.value = guard(item.value)
    for section in plan.sections:
        section.heading = guard(section.heading)
        section.body = guard(section.body)
    if plan.table:
        plan.table.headers = [guard(value) for value in plan.table.headers]
        plan.table.rows = [[guard(value) for value in row] for row in plan.table.rows]
    return plan


def _content(prompt: str, source: str) -> tuple[str, list[tuple[str, str]]]:
    title = _clean_title(prompt)
    source_sentences = [part.strip() for part in re.split(r"(?<=[.!?])\s+", source) if len(part.strip()) > 35]
    seed = source_sentences[:8]
    sections = [
        ("Executive summary", seed[0] if seed else f"This document presents a focused overview of {title.lower()}, the key considerations, and a practical path forward."),
        ("Current context", seed[1] if len(seed) > 1 else f"The opportunity is to turn {title.lower()} into a clear, audience-ready narrative supported by concrete priorities."),
        ("Key insights", seed[2] if len(seed) > 2 else "The strongest approach combines a concise message, evidence-led decisions, and a structure that makes the next action obvious."),
        ("Recommended direction", seed[3] if len(seed) > 3 else "Prioritize the highest-value outcome, define ownership, and measure progress through a small set of meaningful indicators."),
        ("Next steps", seed[4] if len(seed) > 4 else "Align stakeholders, confirm the first milestone, and begin with a focused implementation that can be reviewed and improved quickly."),
    ]
    return title, sections


async def _ai_content(payload: CreateRequest, source: str) -> GeneratedContent:
    settings = get_settings()
    if not settings.llm_api_key:
        raise HTTPException(status_code=503, detail="InsightPDF document generation is not configured")
    format_name = {"docx": "Word document", "pdf": "PDF report", "pptx": "PowerPoint presentation"}[payload.output_format]
    source_rule = (
        "Use the supplied source text as evidence. Do not invent figures, names, dates, or claims that are absent from it."
        if source.strip()
        else "Develop useful, specific content from the user's brief. Avoid fabricated citations, quotations, and unsupported statistics."
    )
    system = (
        f"You are a senior document strategist, editor, and information designer creating a production-ready {format_name}. {source_rule} "
        "Return only one JSON object with this exact shape: "
        '{"title":"...","subtitle":"...","document_type":"invoice|proposal|report|policy|resume|letter|agenda|presentation|general",'
        '"layout":"business|editorial|compact|formal|modern","accent_color":"#RRGGBB",'
        '"metadata":[{"label":"...","value":"..."}],"sections":[{"heading":"...","body":"..."}],'
        '"table":{"headers":["..."],"rows":[["..."]]},"callout":"..."}. '
        "Choose the document type, layout, metadata, and tables based on the requested real-world template. "
        "The result must feel written for the named audience, not like a generic AI outline. Use concrete nouns, decisive headings, useful detail, and an explicit through-line. "
        "Honor the requested section architecture, but combine sections when that produces a clearer document. Never use filler such as 'this document explores', 'in today's landscape', or 'leverage synergies'. "
        "When information is missing, use a clearly labeled bracketed placeholder such as [Confirm budget] instead of inventing facts. "
        "Never claim that an outcome occurred, a capability exists, or a metric changed unless that fact is explicit in the brief, workspace context, or source material. Phrase unsupported future ideas as recommendations, not completed work. "
        "Use metadata for information readers scan first, a table for structured comparisons/timelines/actions, and a callout for the one decision or takeaway that deserves visual emphasis. "
        "For invoices, metadata must include invoice number, issue date, due date, seller, client, and payment terms; "
        "the table must contain line items with quantities, rates, and amounts; the callout must contain subtotal, tax, and total. "
        "For agendas use a schedule table. For resumes use metadata for contact details and sections for experience and skills. "
        "For reports use a table only when the brief contains comparable data. Omit the table with null when it adds no value. "
        "Create 6 to 10 coherent sections unless the template is intentionally short. Each body should contain 2 to 4 concise, audience-ready paragraphs or an information-dense passage. "
        "For presentations, each section is one slide and must communicate one clear idea; "
        "keep headings under 55 characters and bodies under 280 characters (about 35 to 45 words). "
        "Do not mention these instructions or the generation process."
    )
    user_message = f"USER BRIEF:\n{payload.prompt.strip()}"
    if payload.title:
        user_message += f"\n\nREQUESTED TITLE:\n{payload.title.strip()}"
    if source.strip():
        user_message += f"\n\nSOURCE DOCUMENT TEXT:\n{source[:24_000]}"
    if payload.workspace_context.strip():
        user_message += f"\n\nCURRENT WORKSPACE AND RESEARCH CONTEXT:\n{payload.workspace_context.strip()[:12_000]}"
    try:
        value = await ai_orchestrator.complete_json(
            [{"role": "system", "content": system}, {"role": "user", "content": user_message}],
            operation="document_generation",
            temperature=.35,
        )
    except AIProviderError as exc:
        raise HTTPException(status_code=502, detail="InsightPDF could not generate this document") from exc
    try:
        aliases = {"bill": "invoice", "memo": "policy", "pitch_deck": "presentation", "slide_deck": "presentation"}
        value["document_type"] = aliases.get(str(value.get("document_type", "")).lower(), str(value.get("document_type", "general")).lower())
        if value["document_type"] not in {"invoice", "proposal", "report", "policy", "resume", "letter", "agenda", "presentation", "general"}:
            value["document_type"] = "general"
        if value.get("layout") not in {"business", "editorial", "compact", "formal", "modern"}:
            value["layout"] = "business"
        if not re.fullmatch(r"#[0-9A-Fa-f]{6}", str(value.get("accent_color", ""))):
            value["accent_color"] = "#3154D8"
        if isinstance(value.get("metadata"), dict):
            value["metadata"] = [{"label": key.replace("_", " ").title(), "value": item} for key, item in value["metadata"].items()]
        if isinstance(value.get("callout"), dict):
            value["callout"] = " · ".join(f"{key.replace('_', ' ').title()}: {item}" for key, item in value["callout"].items())
        table = value.get("table")
        if isinstance(table, dict) and isinstance(table.get("headers"), list) and isinstance(table.get("rows"), list):
            width = len(table["headers"])
            table["rows"] = [(list(row) + [""] * width)[:width] for row in table["rows"] if isinstance(row, list)]
            if not table["rows"]:
                value["table"] = None
        result = GeneratedContent.model_validate(value)
    except ValueError as exc:
        raise HTTPException(status_code=502, detail="InsightPDF returned an invalid document structure") from exc
    return result


def _docx(title: str, sections: list[tuple[str, str]], theme: dict[str, str]) -> bytes:
    document = WordDocument()
    section = document.sections[0]
    section.top_margin = Inches(.72)
    section.bottom_margin = Inches(.72)
    section.left_margin = Inches(.82)
    section.right_margin = Inches(.82)
    heading = document.add_heading(title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    heading.runs[0].font.color.rgb = RGBColor.from_string(theme["navy"])
    heading.runs[0].font.size = Pt(28)
    sub = document.add_paragraph("Created with InsightPDF AI")
    sub.runs[0].font.color.rgb = RGBColor.from_string(theme["accent"])
    sub.runs[0].font.size = Pt(10)
    for index, (name, body) in enumerate(sections, 1):
        heading = document.add_heading(f"{index:02d}  {name}", level=1)
        heading.runs[0].font.color.rgb = RGBColor.from_string(theme["navy"])
        paragraph = document.add_paragraph(body)
        paragraph.paragraph_format.space_after = Pt(12)
        paragraph.paragraph_format.line_spacing = 1.18
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pdf(title: str, sections: list[tuple[str, str]], theme: dict[str, str]) -> bytes:
    pdf = fitz.open()
    page = pdf.new_page()
    accent = tuple(int(theme["accent"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    navy = tuple(int(theme["navy"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    page.draw_rect(fitz.Rect(0, 0, 14, page.rect.height), color=accent, fill=accent)
    page.insert_textbox(fitz.Rect(48, 58, 545, 150), title, fontsize=26, fontname="hebo", color=navy)
    page.insert_text((48, 158), "Created with InsightPDF AI", fontsize=9, color=accent)
    y = 205
    for index, (name, body) in enumerate(sections, 1):
        page.insert_text((48, y), f"{index:02d}  {name.upper()}", fontsize=11, fontname="hebo", color=navy)
        y += 18
        used = page.insert_textbox(fitz.Rect(48, y, 545, y + 72), body, fontsize=10, lineheight=1.25, color=(.26, .3, .38))
        y += max(60, 78 - min(used, 0))
        if y > 730 and index < len(sections):
            page = pdf.new_page()
            y = 60
    return pdf.tobytes(garbage=4, deflate=True)


def _pptx(title: str, sections: list[tuple[str, str]], theme: dict[str, str]) -> bytes:
    deck = Presentation()
    deck.slide_width = PptxInches(13.333)
    deck.slide_height = PptxInches(7.5)
    colors = {key: PptxRGBColor.from_string(value) for key, value in theme.items()}

    def add_text(slide, text, left, top, width, height, size, color, bold=False):
        box = slide.shapes.add_textbox(PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = "Aptos"
        paragraph.font.size = PptxPt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        return box

    title_slide = deck.slides.add_slide(deck.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = colors["navy"]
    add_text(title_slide, "INSIGHTPDF PRESENTATION", .85, .75, 5, .35, 13, colors["accent"], True)
    add_text(title_slide, title, .85, 2.0, 9.4, 2.0, 34, PptxRGBColor(255, 255, 255), True)
    add_text(title_slide, "A clear, AI-created narrative", .85, 5.75, 6, .4, 18, PptxRGBColor(190, 200, 218))
    for index, (name, body) in enumerate(sections, 1):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptxRGBColor(255, 255, 255)
        accent = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(.16), deck.slide_height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = colors["accent"]
        accent.line.fill.background()
        add_text(slide, f"{index:02d} · {name.upper()}", .82, .72, 5, .35, 13, colors["accent"], True)
        add_text(slide, name, .82, 1.45, 10.7, .85, 30, colors["navy"], True)
        add_text(slide, body, .82, 2.65, 8.9, 2.2, 20, PptxRGBColor(62, 72, 91))
        add_text(slide, "InsightPDF", .82, 6.7, 2, .3, 10, PptxRGBColor(130, 140, 158), True)
        add_text(slide, f"{index + 1:02d}", 11.8, 6.7, .5, .3, 10, PptxRGBColor(130, 140, 158))
    output = BytesIO()
    deck.save(output)
    return output.getvalue()


def _shade_cell(cell, color: str) -> None:
    properties = cell._tc.get_or_add_tcPr()
    fill = OxmlElement("w:shd")
    fill.set(qn("w:fill"), color)
    properties.append(fill)


def _docx_dynamic(plan: GeneratedContent, theme: dict[str, str]) -> bytes:
    document = WordDocument()
    section = document.sections[0]
    section.top_margin = Inches(.68)
    section.bottom_margin = Inches(.68)
    section.left_margin = Inches(.78)
    section.right_margin = Inches(.78)
    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10.5 if plan.layout == "compact" else 11)
    for style_name, size in (("Heading 1", 16), ("Heading 2", 13)):
        styles[style_name].font.name = "Aptos Display"
        styles[style_name].font.size = Pt(size)
        styles[style_name].font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))

    if plan.document_type == "invoice":
        brand = document.add_paragraph()
        brand.paragraph_format.space_after = Pt(3)
        run = brand.add_run(next((item.value for item in plan.metadata if item.label.lower() in {"seller", "from"}), "INVOICE"))
        run.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor.from_string(theme["navy"])
        title = document.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        title.paragraph_format.space_after = Pt(16)
        title_run = title.add_run("INVOICE")
        title_run.bold = True
        title_run.font.size = Pt(30)
        title_run.font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))
        details = document.add_table(rows=0, cols=2)
        details.alignment = WD_TABLE_ALIGNMENT.CENTER
        details.autofit = False
        for item in plan.metadata:
            cells = details.add_row().cells
            cells[0].width, cells[1].width = Inches(1.55), Inches(4.95)
            cells[0].text, cells[1].text = item.label.upper(), item.value
            cells[0].paragraphs[0].runs[0].font.bold = True
            cells[0].paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))
        document.add_paragraph()
    else:
        kicker = document.add_paragraph(plan.document_type.upper())
        kicker.paragraph_format.space_after = Pt(6)
        kicker.runs[0].bold = True
        kicker.runs[0].font.size = Pt(9)
        kicker.runs[0].font.color.rgb = RGBColor.from_string(plan.accent_color.removeprefix("#"))
        heading = document.add_paragraph()
        heading.paragraph_format.space_after = Pt(5)
        heading_run = heading.add_run(plan.title)
        heading_run.bold = True
        heading_run.font.size = Pt(28 if plan.layout == "editorial" else 24)
        heading_run.font.color.rgb = RGBColor.from_string(theme["navy"])
        if plan.subtitle:
            subtitle = document.add_paragraph(plan.subtitle)
            subtitle.paragraph_format.space_after = Pt(15)
            subtitle.runs[0].italic = True
            subtitle.runs[0].font.color.rgb = RGBColor(100, 110, 126)
        if plan.metadata:
            metadata = document.add_table(rows=0, cols=2)
            metadata.alignment = WD_TABLE_ALIGNMENT.CENTER
            for index in range(0, len(plan.metadata), 2):
                cells = metadata.add_row().cells
                for offset, cell in enumerate(cells):
                    if index + offset < len(plan.metadata):
                        item = plan.metadata[index + offset]
                        cell.text = f"{item.label.upper()}\n{item.value}"
                        cell.paragraphs[0].runs[0].font.bold = True
                        _shade_cell(cell, theme["soft"])
            document.add_paragraph()

    if plan.table:
        table = document.add_table(rows=1, cols=len(plan.table.headers))
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.style = "Table Grid"
        table.autofit = True
        for index, header in enumerate(plan.table.headers):
            cell = table.rows[0].cells[index]
            cell.text = header
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            _shade_cell(cell, plan.accent_color.removeprefix("#"))
            for run in cell.paragraphs[0].runs:
                run.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
        for row in plan.table.rows:
            cells = table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = value
                cells[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                if len(table.rows) % 2 == 1:
                    _shade_cell(cells[index], theme["soft"])
        document.add_paragraph()

    if plan.callout:
        callout = document.add_table(rows=1, cols=1)
        callout.alignment = WD_TABLE_ALIGNMENT.RIGHT if plan.document_type == "invoice" else WD_TABLE_ALIGNMENT.CENTER
        callout.cell(0, 0).text = plan.callout
        _shade_cell(callout.cell(0, 0), theme["soft"])
        callout.cell(0, 0).paragraphs[0].runs[0].font.bold = True
        callout.cell(0, 0).paragraphs[0].runs[0].font.color.rgb = RGBColor.from_string(theme["navy"])
        document.add_paragraph()

    visible_sections = plan.sections[:2] if plan.document_type == "invoice" else plan.sections
    for item in visible_sections:
        heading = document.add_heading(item.heading, level=1)
        heading.paragraph_format.keep_with_next = True
        body = document.add_paragraph(item.body)
        body.paragraph_format.space_after = Pt(9)
        body.paragraph_format.line_spacing = 1.15 if plan.layout != "editorial" else 1.28
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.add_run(f"Created with InsightPDF AI · {plan.document_type.title()} · {plan.layout.title()} layout").font.size = Pt(8)
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pdf_dynamic(plan: GeneratedContent, theme: dict[str, str]) -> bytes:
    document = fitz.open()
    accent = tuple(int(plan.accent_color[index:index + 2], 16) / 255 for index in (1, 3, 5))
    navy = tuple(int(theme["navy"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    soft = tuple(int(theme["soft"][index:index + 2], 16) / 255 for index in (0, 2, 4))
    if plan.document_type != "invoice":
        muted = (.35, .39, .47)
        light_ink = (.77, .8, .86)

        def footer(target, number: int) -> None:
            target.draw_line(fitz.Point(48, 795), fitz.Point(545, 795), color=(.86, .87, .9), width=.5)
            target.insert_text((48, 815), plan.title[:70], fontsize=7, color=(.55, .58, .64))
            target.insert_text((526, 815), f"{number:02d}", fontsize=7, fontname="hebo", color=accent)

        cover = document.new_page()
        cover.draw_rect(fitz.Rect(0, 0, cover.rect.width, 245), color=navy, fill=navy)
        cover.draw_rect(fitz.Rect(0, 0, 15, cover.rect.height), color=accent, fill=accent)
        cover.insert_text((48, 55), plan.document_type.upper(), fontsize=9, fontname="hebo", color=accent)
        cover.insert_textbox(fitz.Rect(48, 82, 545, 188), plan.title, fontsize=29, fontname="hebo", lineheight=1.05, color=(1, 1, 1))
        if plan.subtitle:
            cover.insert_textbox(fitz.Rect(48, 190, 520, 232), plan.subtitle, fontsize=10, lineheight=1.25, color=light_ink)
        cover.insert_text((48, 284), "DOCUMENT BRIEF", fontsize=7, fontname="hebo", color=accent)
        y = 306
        for index, item in enumerate(plan.metadata[:6]):
            column, row = index % 2, index // 2
            x, py = 48 + column * 250, y + row * 56
            cover.draw_rect(fitz.Rect(x, py, x + 232, py + 45), color=soft, fill=soft)
            cover.insert_text((x + 10, py + 14), item.label.upper(), fontsize=6.5, fontname="hebo", color=accent)
            cover.insert_textbox(fitz.Rect(x + 10, py + 19, x + 222, py + 40), item.value, fontsize=8.5, color=navy)
        y += max(1, (min(len(plan.metadata), 6) + 1) // 2) * 56 + 22
        if plan.callout:
            cover.draw_rect(fitz.Rect(48, y, 545, min(y + 72, 745)), color=soft, fill=soft)
            cover.draw_rect(fitz.Rect(48, y, 53, min(y + 72, 745)), color=accent, fill=accent)
            cover.insert_text((67, y + 20), "KEY TAKEAWAY", fontsize=7, fontname="hebo", color=accent)
            cover.insert_textbox(fitz.Rect(67, y + 27, 525, min(y + 65, 740)), plan.callout, fontsize=9.5, lineheight=1.2, color=navy)
        footer(cover, 1)

        page_number = 2
        if plan.table:
            page = document.new_page()
            page.insert_text((48, 54), "STRUCTURED OVERVIEW", fontsize=8, fontname="hebo", color=accent)
            page.insert_text((48, 88), "At a glance", fontsize=23, fontname="hebo", color=navy)
            table_y = 132
            column_width = 497 / len(plan.table.headers)
            page.draw_rect(fitz.Rect(48, table_y, 545, table_y + 29), color=accent, fill=accent)
            for index, header in enumerate(plan.table.headers):
                page.insert_textbox(fitz.Rect(54 + index * column_width, table_y + 9, 48 + (index + 1) * column_width - 4, table_y + 27), header, fontsize=7.5, fontname="hebo", color=(1, 1, 1))
            table_y += 29
            for row_index, row in enumerate(plan.table.rows[:18]):
                if row_index % 2 == 0:
                    page.draw_rect(fitz.Rect(48, table_y, 545, table_y + 31), color=soft, fill=soft)
                for index, value in enumerate(row):
                    page.insert_textbox(fitz.Rect(54 + index * column_width, table_y + 9, 48 + (index + 1) * column_width - 4, table_y + 29), value, fontsize=7.5, color=navy)
                table_y += 31
            footer(page, page_number)
            page_number += 1

        for section_index in range(0, len(plan.sections), 2):
            page = document.new_page()
            page.draw_rect(fitz.Rect(0, 0, 9, page.rect.height), color=accent, fill=accent)
            page.insert_text((48, 52), f"{section_index + 1:02d} / {len(plan.sections):02d}", fontsize=8, fontname="hebo", color=accent)
            section_y = 92
            for offset, item in enumerate(plan.sections[section_index:section_index + 2]):
                page.insert_text((48, section_y), f"{section_index + offset + 1:02d}", fontsize=10, fontname="hebo", color=accent)
                page.insert_textbox(fitz.Rect(80, section_y - 17, 545, section_y + 32), item.heading, fontsize=20, fontname="hebo", lineheight=1.05, color=navy)
                section_y += 48
                page.insert_textbox(fitz.Rect(80, section_y, 535, section_y + 245), item.body, fontsize=10.5, lineheight=1.35, color=muted)
                if offset == 0 and section_index + 1 < len(plan.sections):
                    page.draw_line(fitz.Point(80, section_y + 265), fitz.Point(535, section_y + 265), color=soft, width=2)
                section_y += 330
            footer(page, page_number)
            page_number += 1
        return document.tobytes(garbage=4, deflate=True)

    page = document.new_page()
    page.draw_rect(fitz.Rect(0, 0, 12, page.rect.height), color=accent, fill=accent)
    if plan.document_type == "invoice":
        page.insert_text((48, 66), "INVOICE", fontsize=30, fontname="hebo", color=accent)
        page.insert_text((48, 92), plan.title, fontsize=11, color=navy)
    else:
        page.insert_text((48, 54), plan.document_type.upper(), fontsize=9, fontname="hebo", color=accent)
        page.insert_textbox(fitz.Rect(48, 72, 545, 145), plan.title, fontsize=25, fontname="hebo", color=navy)
        if plan.subtitle:
            page.insert_textbox(fitz.Rect(48, 140, 545, 180), plan.subtitle, fontsize=10, color=(.38, .42, .5))
    y = 130 if plan.document_type == "invoice" else 190
    if plan.metadata:
        for index, item in enumerate(plan.metadata):
            column = index % 2
            row = index // 2
            x = 48 + column * 250
            py = y + row * 35
            page.insert_text((x, py), item.label.upper(), fontsize=7, fontname="hebo", color=accent)
            page.insert_textbox(fitz.Rect(x, py + 5, x + 225, py + 30), item.value, fontsize=9, color=navy)
        y += ((len(plan.metadata) + 1) // 2) * 35 + 15
    if plan.table:
        column_width = 497 / len(plan.table.headers)
        page.draw_rect(fitz.Rect(48, y, 545, y + 25), color=accent, fill=accent)
        for index, header in enumerate(plan.table.headers):
            page.insert_textbox(fitz.Rect(52 + index * column_width, y + 7, 48 + (index + 1) * column_width - 4, y + 23), header, fontsize=8, fontname="hebo", color=(1, 1, 1))
        y += 25
        for row_index, row in enumerate(plan.table.rows):
            if y > 720:
                page = document.new_page()
                y = 55
            if row_index % 2 == 0:
                page.draw_rect(fitz.Rect(48, y, 545, y + 27), color=soft, fill=soft)
            for index, value in enumerate(row):
                page.insert_textbox(fitz.Rect(52 + index * column_width, y + 7, 48 + (index + 1) * column_width - 4, y + 26), value, fontsize=8, color=navy)
            y += 27
        y += 12
    if plan.callout:
        callout_text = plan.callout.replace(" · ", "\n")
        callout_height = 66 if plan.document_type == "invoice" else 48
        page.draw_rect(fitz.Rect(315, y, 545, y + callout_height), color=soft, fill=soft)
        page.insert_textbox(fitz.Rect(328, y + 9, 533, y + callout_height - 7), callout_text, fontsize=9, fontname="hebo", lineheight=1.15, align=2, color=navy)
        y += callout_height + 18
    if plan.document_type == "invoice":
        for index, item in enumerate(plan.sections[:2]):
            x = 48 + index * 250
            page.insert_text((x, y), item.heading.upper(), fontsize=9, fontname="hebo", color=accent)
            page.insert_textbox(fitz.Rect(x, y + 12, x + 225, min(page.rect.height - 42, y + 105)), item.body, fontsize=8.5, lineheight=1.18, color=(.25, .29, .36))
        return document.tobytes(garbage=4, deflate=True)
    for item in plan.sections:
        if y > 700:
            page = document.new_page()
            y = 58
        page.insert_text((48, y), item.heading.upper(), fontsize=10, fontname="hebo", color=accent)
        y += 15
        page.insert_textbox(fitz.Rect(48, y, 545, y + 74), item.body, fontsize=9.5, lineheight=1.25, color=(.25, .29, .36))
        y += 77
    return document.tobytes(garbage=4, deflate=True)


def _shorten_slide_text(value: str, limit: int) -> str:
    value = re.sub(r"\s+", " ", value).strip()
    if len(value) <= limit:
        return value
    candidate = value[:limit - 1].rsplit(" ", 1)[0].rstrip(" ,;:-")
    return f"{candidate or value[:limit - 1].rstrip()}…"


def _pptx_preview_slides(plan: GeneratedContent) -> list[dict[str, str]]:
    slides = [{
        "eyebrow": _shorten_slide_text(plan.document_type.upper(), 34),
        "title": _shorten_slide_text(plan.title, 76),
        "body": _shorten_slide_text(plan.subtitle or "Created with InsightPDF AI", 130),
        "variant": "title",
    }]
    limits = {1: (64, 330), 2: (48, 250), 3: (68, 310)}
    for index, section in enumerate(plan.sections, 1):
        variant = ((index - 1) % 3) + 1
        heading_limit, body_limit = limits[variant]
        slides.append({
            "eyebrow": f"{index:02d}",
            "title": _shorten_slide_text(section.heading, heading_limit),
            "body": _shorten_slide_text(section.body, body_limit),
            "variant": f"section-{variant}",
        })
    return slides


def _document_preview_pages(plan: GeneratedContent) -> list[dict]:
    """Return the same structured content the renderer used, grouped into preview pages."""
    pages: list[dict] = [{
        "kind": "cover",
        "eyebrow": plan.document_type.upper(),
        "title": plan.title,
        "subtitle": plan.subtitle,
        "metadata": [item.model_dump() for item in plan.metadata[:6]],
        "callout": plan.callout,
    }]
    if plan.table:
        pages.append({
            "kind": "table",
            "eyebrow": "AT A GLANCE",
            "title": next((section.heading for section in plan.sections if any(word in section.heading.lower() for word in ("summary", "overview", "scorecard", "timeline", "action"))), "Structured overview"),
            "table": plan.table.model_dump(),
            "callout": plan.callout,
        })
    for index in range(0, len(plan.sections), 2):
        pages.append({
            "kind": "sections",
            "eyebrow": f"{index + 1:02d} / {len(plan.sections):02d}",
            "sections": [item.model_dump() for item in plan.sections[index:index + 2]],
        })
    return pages


def _native_blocks(plan: GeneratedContent) -> list[dict[str, str]]:
    blocks: list[dict[str, str]] = []
    if plan.subtitle:
        blocks.append({"type": "paragraph", "text": plan.subtitle})
    for item in plan.metadata:
        blocks.append({"type": "bullet", "text": f"{item.label}: {item.value}"})
    if plan.callout:
        blocks.extend([
            {"type": "heading", "text": "Key takeaway"},
            {"type": "paragraph", "text": plan.callout},
        ])
    if plan.table:
        blocks.append({"type": "heading", "text": "Structured overview"})
        blocks.append({"type": "paragraph", "text": " | ".join(plan.table.headers)})
        blocks.extend({"type": "bullet", "text": " | ".join(row)} for row in plan.table.rows)
    for item in plan.sections:
        blocks.extend([
            {"type": "heading", "text": item.heading},
            {"type": "paragraph", "text": item.body},
        ])
    return blocks or [{"type": "paragraph", "text": ""}]


def _pptx_dynamic(plan: GeneratedContent, theme: dict[str, str], hero_image: bytes | None = None) -> bytes:
    slide_data = _pptx_preview_slides(plan)
    sections = [(item["title"], item["body"]) for item in slide_data[1:]]
    deck = Presentation()
    deck.slide_width = PptxInches(13.333)
    deck.slide_height = PptxInches(7.5)
    navy = PptxRGBColor.from_string(theme["navy"])
    accent = PptxRGBColor.from_string(plan.accent_color.removeprefix("#"))
    soft = PptxRGBColor.from_string(theme["soft"])

    def text(slide, value, left, top, width, height, size, color, bold=False,
             align=PP_ALIGN.LEFT, vertical=MSO_ANCHOR.TOP):
        shape = slide.shapes.add_textbox(PptxInches(left), PptxInches(top), PptxInches(width), PptxInches(height))
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = vertical
        paragraph = frame.paragraphs[0]
        paragraph.text = value
        paragraph.alignment = align
        paragraph.space_before = paragraph.space_after = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = PptxPt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        return shape

    title_slide = deck.slides.add_slide(deck.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = navy
    title_data = slide_data[0]
    if hero_image:
        try:
            title_slide.shapes.add_picture(BytesIO(hero_image), PptxInches(8.7), PptxInches(0), width=PptxInches(4.633), height=PptxInches(7.5))
        except Exception:
            hero_image = None
    text(title_slide, title_data["eyebrow"], .85, .72, 5, .3, 14, accent, True)
    text(title_slide, title_data["title"], .85, 1.65, 7.35 if hero_image else 11.25, 2.25, 50, PptxRGBColor(255, 255, 255), True,
         vertical=MSO_ANCHOR.MIDDLE)
    text(title_slide, title_data["body"], .85, 5.55, 7.15 if hero_image else 9.5, .85, 19, PptxRGBColor(190, 200, 218))
    for index, (heading, body) in enumerate(sections, 1):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PptxRGBColor(255, 255, 255) if index % 3 else navy
        ink = navy if index % 3 else PptxRGBColor(255, 255, 255)
        muted = PptxRGBColor(62, 72, 91) if index % 3 else PptxRGBColor(195, 204, 219)
        if index % 3 == 1:
            text(slide, f"{index:02d}", .8, .65, 1, .35, 18, accent, True)
            text(slide, heading, .8, 1.35, 9.25, 1.15, 35, ink, True, vertical=MSO_ANCHOR.MIDDLE)
            text(slide, body, .8, 2.95, 8.95, 2.55, 19, muted)
            marker = slide.shapes.add_shape(1, PptxInches(10.75), PptxInches(1.35), PptxInches(1.65), PptxInches(4.85))
            marker.fill.solid()
            marker.fill.fore_color.rgb = accent
            marker.line.fill.background()
        elif index % 3 == 2:
            panel = slide.shapes.add_shape(1, PptxInches(0), PptxInches(0), PptxInches(4.9), deck.slide_height)
            panel.fill.solid()
            panel.fill.fore_color.rgb = soft
            panel.line.fill.background()
            text(slide, f"{index:02d}", 1.05, 2.35, 2.7, 1.4, 60, accent, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            text(slide, heading, 5.75, 1.2, 6.55, 1.45, 35, ink, True, vertical=MSO_ANCHOR.MIDDLE)
            text(slide, body, 5.75, 3.0, 6.35, 2.4, 18, muted)
        else:
            text(slide, f"{index:02d}  /  INSIGHT", .85, .78, 4, .3, 14, accent, True)
            text(slide, heading, .85, 1.75, 11.45, 1.35, 37, ink, True, vertical=MSO_ANCHOR.MIDDLE)
            text(slide, body, .85, 3.5, 10.9, 2.25, 19, muted)
        text(slide, "InsightPDF", .8, 6.8, 2, .25, 10, muted, True)
    output = BytesIO()
    deck.save(output)
    return output.getvalue()


@router.get("/themes")
async def themes() -> list[dict[str, str]]:
    return [
        {"id": "minimal", "name": "Minimal", "description": "Crisp navy and red"},
        {"id": "executive", "name": "Executive", "description": "Charcoal and gold"},
        {"id": "modern", "name": "Modern", "description": "Bold blue and white"},
        {"id": "warm", "name": "Warm", "description": "Editorial earth tones"},
    ]


@router.post("", response_model=ArtifactResponse, status_code=status.HTTP_201_CREATED)
async def create_file(
    payload: CreateRequest,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> GeneratedArtifact:
    source_ids = list(dict.fromkeys([
        *payload.source_document_ids,
        *([payload.source_document_id] if payload.source_document_id else []),
    ]))[:12]
    source_parts: list[str] = []
    for source_id in source_ids:
        document = await owned_document(source_id, user, session)
        pages = list(await session.scalars(
            select(DocumentPage).where(DocumentPage.document_id == source_id).order_by(DocumentPage.page_number)
        ))
        source_parts.append(f"[SOURCE: {document.display_title or document.filename}]\n" + " ".join(page.text for page in pages))
    source_text = "\n\n".join(source_parts)[:24_000]
    from app.deliverables import ensure_personal_workspace
    workspace = await ensure_personal_workspace(user, session)
    server_context: list[str] = []
    if payload.conversation_id:
        from app.chat import owned_conversation
        conversation = await owned_conversation(payload.conversation_id, user, session)
        server_context.append(f"Active research thread: {conversation.title}")
        server_context.extend(
            f"{message.role.value}: {message.content[:1200]}"
            for message in conversation.messages[-8:]
        )
    from app.models import WorkspaceMemory
    memories = list(await session.scalars(
        select(WorkspaceMemory).where(
            WorkspaceMemory.workspace_id == workspace.id,
            WorkspaceMemory.owner_id == user.id,
        ).order_by(WorkspaceMemory.updated_at.desc()).limit(20)
    ))
    if memories:
        server_context.append("Workspace memory:\n" + "\n".join(f"- {item.key}: {item.value}" for item in memories))
    recent_deliverables = list(await session.scalars(
        select(NativeDocument).where(
            NativeDocument.workspace_id == workspace.id,
            NativeDocument.owner_id == user.id,
        ).order_by(NativeDocument.updated_at.desc()).limit(8)
    ))
    if recent_deliverables:
        server_context.append("Recent deliverables: " + ", ".join(item.title for item in recent_deliverables))
    effective_context = "\n\n".join([payload.workspace_context.strip(), *server_context]).strip()[:12000]
    payload = payload.model_copy(update={"workspace_context": effective_context})
    template = CREATION_TEMPLATES.get(payload.template_id or "")
    if payload.template_id and template is None:
        raise HTTPException(status_code=422, detail="Unknown document template")
    content_request = payload
    if template:
        answered_context = "\n".join(f"- {label}: {value}" for label, value in payload.template_answers.items())
        template_prompt = (
            f"TEMPLATE ARCHITECTURE:\n{template['brief']}\n\n"
            f"INTAKE ANSWERS:\n{answered_context or 'No additional intake answers supplied.'}\n\n"
            f"TOPIC AND REQUIREMENTS:\n{payload.prompt}"
        )
        content_request = payload.model_copy(update={"prompt": template_prompt})
    plan = await _ai_content(content_request, source_text)
    grounding_evidence = "\n".join([
        payload.prompt,
        payload.workspace_context,
        source_text,
        *payload.template_answers.values(),
    ])
    plan = _guard_unsupported_metrics(plan, grounding_evidence)
    await record_ai_usage(user, f"create_{payload.output_format}", session, cached=False)
    selected_theme = template["theme"] if template else payload.theme
    if template:
        plan.layout = template["layout"]
        if payload.output_format != "pptx":
            plan.document_type = template["type"]
    theme = THEMES[selected_theme]
    data = _pptx_dynamic(plan, theme) if payload.output_format == "pptx" else {"docx": _docx_dynamic, "pdf": _pdf_dynamic}[payload.output_format](plan, theme)
    filename = safe_filename(re.sub(r"[^a-zA-Z0-9 -]", "", plan.title).strip().replace(" ", "-").lower()[:70] or "insight-document")
    filename = f"{filename}.{payload.output_format}"
    content_types = {
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pdf": "application/pdf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }
    identifier = uuid.uuid4()
    native_identifier = uuid.uuid4()
    key = f"{user.id}/generated/{identifier}/{filename}"
    try:
        ObjectStorage().upload(key, data, content_types[payload.output_format])
    except Exception as exc:
        raise HTTPException(status_code=503, detail="Generated-file storage is temporarily unavailable") from exc
    native_content = {"type": "doc", "blocks": _native_blocks(plan)}
    native_document = NativeDocument(
        id=native_identifier,
        workspace_id=workspace.id,
        owner_id=user.id,
        title=plan.title,
        content=native_content,
    )
    session.add(native_document)
    session.add(NativeDocumentVersion(
        native_document_id=native_identifier,
        version_number=1,
        title=plan.title,
        content=native_content,
        change_summary=f"Generated from {payload.template_id or 'custom brief'}",
        created_by=user.id,
    ))
    for source_id in source_ids:
        session.add(NativeDocumentSource(native_document_id=native_identifier, document_id=source_id))
    artifact = GeneratedArtifact(
        id=identifier,
        owner_id=user.id,
        workspace_id=workspace.id,
        operation=f"ai_create_{payload.output_format}",
        filename=filename,
        object_key=key,
        content_type=content_types[payload.output_format],
        size_bytes=len(data),
        parameters={
            "prompt": payload.prompt,
            "theme": selected_theme,
            "template_id": payload.template_id,
            "source_document_id": str(source_ids[0]) if source_ids else None,
            "source_document_ids": [str(source_id) for source_id in source_ids],
            "source_count": len(source_ids),
            "research_context_used": bool(payload.workspace_context.strip()),
            "native_document_id": str(native_identifier),
            "provider": "mistral",
            "model": get_settings().llm_model,
            "document_type": plan.document_type,
            "layout": plan.layout,
            "accent_color": plan.accent_color,
            "preview_palette": {
                "navy": f"#{theme['navy']}",
                "accent": plan.accent_color,
                "soft": f"#{theme['soft']}",
            },
            "preview_slides": _pptx_preview_slides(plan) if payload.output_format == "pptx" else [],
            "preview_pages": _document_preview_pages(plan) if payload.output_format != "pptx" else [],
        },
    )
    session.add(artifact)
    session.add(ArtifactVersion(
        artifact_id=identifier,
        version_number=1,
        object_key=key,
        content_type=content_types[payload.output_format],
        size_bytes=len(data),
        change_prompt=payload.prompt,
        metadata_json={"operation": artifact.operation, "theme": selected_theme},
    ))
    await session.commit()
    await session.refresh(artifact)
    return artifact


@router.post("/jobs", response_model=ProcessingJobResponse, status_code=status.HTTP_202_ACCEPTED)
async def create_file_job(
    payload: CreateRequest,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
):
    """Queue generation so it survives navigation and exposes progress/cancellation."""
    source_ids = list(dict.fromkeys([
        *payload.source_document_ids,
        *([payload.source_document_id] if payload.source_document_id else []),
    ]))[:12]
    for source_id in source_ids:
        await owned_document(source_id, user, session)
    if payload.template_id and payload.template_id not in CREATION_TEMPLATES:
        raise HTTPException(status_code=422, detail="Unknown document template")
    from app.jobs import create_job_without_documents
    return await create_job_without_documents(
        "ai_create",
        {"request": payload.model_dump(mode="json")},
        user,
        session,
    )
