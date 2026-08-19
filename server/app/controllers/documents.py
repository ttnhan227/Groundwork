import hashlib
import re
import tempfile
import textwrap
import uuid
import zipfile
from datetime import UTC, datetime
from io import BytesIO

import fitz
from fastapi import (
    APIRouter,
    Depends,
    File,
    Form,
    HTTPException,
    Response,
    UploadFile,
    status,
)
from fastapi.responses import StreamingResponse
from PIL import Image, ImageOps, UnidentifiedImageError
from pptx import Presentation
from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession
from starlette.concurrency import run_in_threadpool

from app.config import get_settings
from app.database import get_session
from app.dependencies import current_user
from app.document_conversions import docx_to_markdown, validate_docx
from app.models import (
    Document,
    DocumentPage,
    DocumentStatus,
    GeneratedArtifact,
    JobStatus,
    ProcessingJob,
    User,
)
from app.schemas import (
    DocumentArchiveRequest,
    DocumentPageResponse,
    DocumentRenameRequest,
    DocumentResponse,
    ProcessingJobResponse,
)
from app.storage import ObjectStorage

router = APIRouter(prefix="/documents", tags=["Documents"])


def safe_filename(name: str) -> str:
    basename = name.replace("\\", "/").split("/")[-1]
    value = re.sub(r"[^a-zA-Z0-9._ -]", "_", basename).strip(" .")
    return value[:180] or "document.pdf"


IMAGE_CONTENT_TYPES = {"image/jpeg", "image/png", "image/webp"}
IMAGE_FORMATS = {"JPEG", "PNG", "WEBP"}
MAX_IMAGE_PIXELS = 40_000_000
TEXT_CONTENT_TYPES = {"text/plain", "text/markdown", "text/rtf", "application/rtf"}
DOCX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def text_to_pdf(text: str, title: str) -> bytes:
    """Create a readable normalized derivative without claiming native-format fidelity."""
    cleaned = text.replace("\x00", "").strip()
    if not cleaned:
        raise ValueError("The document does not contain readable text")
    pdf = fitz.open()
    try:
        lines: list[str] = []
        for raw_line in cleaned.splitlines():
            if not raw_line.strip():
                lines.append("")
                continue
            lines.extend(textwrap.wrap(
                raw_line,
                width=88,
                replace_whitespace=False,
                drop_whitespace=True,
                break_long_words=True,
                break_on_hyphens=False,
            ) or [""])
        lines_per_page = 42
        for offset in range(0, len(lines), lines_per_page):
            page = pdf.new_page(width=595, height=842)
            if offset == 0:
                page.insert_textbox(fitz.Rect(52, 42, 543, 78), title, fontsize=16, fontname="helv", color=(0.08, 0.12, 0.24))
            top = 92 if offset == 0 else 52
            page.insert_textbox(
                fitz.Rect(52, top, 543, 790),
                "\n".join(lines[offset:offset + lines_per_page]),
                fontsize=10,
                fontname="helv",
                lineheight=1.35,
                color=(0.12, 0.15, 0.22),
            )
        return pdf.tobytes(garbage=4, deflate=True)
    finally:
        pdf.close()


def source_to_pdf(data: bytes, suffix: str, title: str) -> bytes:
    if suffix == "pdf":
        if not data.startswith(b"%PDF-"):
            raise ValueError("The uploaded file is not a valid PDF")
        return data
    if suffix == "docx":
        validate_docx(data)
        return text_to_pdf(docx_to_markdown(data).decode("utf-8"), title)
    if suffix == "pptx":
        try:
            presentation = Presentation(BytesIO(data))
        except Exception as exc:
            raise ValueError("The uploaded file is not a valid PPTX presentation") from exc
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            values = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            slides.append(f"Slide {index}\n" + "\n".join(values))
        return text_to_pdf("\n\n".join(slides), title)
    try:
        text = data.decode("utf-8-sig")
    except UnicodeDecodeError:
        try:
            text = data.decode("cp1252")
        except UnicodeDecodeError as exc:
            raise ValueError("The text encoding is not supported") from exc
    if suffix == "rtf":
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?|[{}]", "", text)
    return text_to_pdf(text, title)


def image_to_pdf(data: bytes) -> bytes:
    """Validate and normalize a supported image into a one-page PDF."""
    try:
        with Image.open(BytesIO(data)) as source:
            if source.format not in IMAGE_FORMATS:
                raise ValueError("Only PNG, JPEG, and WebP images are accepted")
            if source.width * source.height > MAX_IMAGE_PIXELS:
                raise ValueError("Image dimensions are too large")
            source.load()
            image = ImageOps.exif_transpose(source)
            if image.mode in {"RGBA", "LA"} or "transparency" in image.info:
                rgba = image.convert("RGBA")
                normalized = Image.new("RGB", rgba.size, "white")
                normalized.paste(rgba, mask=rgba.getchannel("A"))
            else:
                normalized = image.convert("RGB")
            output = BytesIO()
            normalized.save(output, format="PDF", resolution=150.0, quality=92)
            return output.getvalue()
    except (UnidentifiedImageError, OSError) as exc:
        raise ValueError("The file is not a valid supported image") from exc


@router.get("", response_model=list[DocumentResponse])
async def list_documents(
    workspace_id: uuid.UUID | None = None,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> list[Document]:
    query = select(Document).where(Document.owner_id == user.id)
    if workspace_id is not None:
        query = query.where(Document.workspace_id == workspace_id)
    result = await session.scalars(query.order_by(Document.created_at.desc()))
    return list(result)


def unique_archive_name(filename: str, used_names: set[str]) -> str:
    safe_name = safe_filename(filename)
    stem, separator, suffix = safe_name.rpartition(".")
    if not separator:
        stem, suffix = safe_name, ""
    candidate = safe_name
    counter = 2
    while candidate.casefold() in used_names:
        candidate = f"{stem} ({counter}){'.' + suffix if suffix else ''}"
        counter += 1
    used_names.add(candidate.casefold())
    return candidate


def build_documents_archive(files: list[tuple[str, str]]):
    archive = tempfile.SpooledTemporaryFile(max_size=1024 * 1024, mode="w+b")
    storage = ObjectStorage()
    used_names: set[str] = set()
    try:
        with zipfile.ZipFile(archive, mode="w", compression=zipfile.ZIP_STORED) as bundle:
            for filename, object_key in files:
                bundle.writestr(
                    unique_archive_name(filename, used_names),
                    storage.download(object_key),
                )
        archive.seek(0)
        return archive
    except Exception:
        archive.close()
        raise


@router.post("/download-zip")
async def download_documents_archive(
    payload: DocumentArchiveRequest,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> StreamingResponse:
    references = list(dict.fromkeys((item.kind, item.id) for item in payload.files))
    if len(references) < 2:
        raise HTTPException(status_code=422, detail="Select at least two different files")
    document_ids = [identifier for kind, identifier in references if kind == "document"]
    artifact_ids = [identifier for kind, identifier in references if kind == "artifact"]
    documents = list(await session.scalars(
        select(Document).where(Document.id.in_(document_ids), Document.owner_id == user.id)
    )) if document_ids else []
    artifacts = list(await session.scalars(
        select(GeneratedArtifact).where(GeneratedArtifact.id.in_(artifact_ids), GeneratedArtifact.owner_id == user.id)
    )) if artifact_ids else []
    document_map = {item.id: item for item in documents}
    artifact_map = {item.id: item for item in artifacts}
    if len(document_map) != len(document_ids) or len(artifact_map) != len(artifact_ids):
        raise HTTPException(status_code=404, detail="One or more files were not found")
    files = [
        (document_map[identifier].filename, document_map[identifier].object_key)
        if kind == "document"
        else (artifact_map[identifier].filename, artifact_map[identifier].object_key)
        for kind, identifier in references
    ]
    total_size = sum(document_map[item].size_bytes for item in document_ids) + sum(artifact_map[item].size_bytes for item in artifact_ids)
    if total_size > 500 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="Selected documents exceed the 500 MB archive limit")

    try:
        archive = await run_in_threadpool(build_documents_archive, files)
    except Exception:
        raise HTTPException(status_code=503, detail="Could not prepare the document archive")

    def stream_archive():
        try:
            while chunk := archive.read(1024 * 1024):
                yield chunk
        finally:
            archive.close()

    return StreamingResponse(
        stream_archive(),
        media_type="application/zip",
        headers={
            "Content-Disposition": 'attachment; filename="groundwork-documents.zip"',
            "X-File-Count": str(len(files)),
        },
    )


async def owned_document(document_id: uuid.UUID, user: User, session: AsyncSession) -> Document:
    document = await session.scalar(select(Document).where(Document.id == document_id, Document.owner_id == user.id))
    if document is None:
        raise HTTPException(status_code=404, detail="Document not found")
    return document


@router.get("/{document_id}", response_model=DocumentResponse)
async def get_document(document_id: uuid.UUID, user: User = Depends(current_user), session: AsyncSession = Depends(get_session)) -> Document:
    return await owned_document(document_id, user, session)


@router.patch("/{document_id}", response_model=DocumentResponse)
async def rename_document(
    document_id: uuid.UUID, payload: DocumentRenameRequest,
    user: User = Depends(current_user), session: AsyncSession = Depends(get_session),
) -> Document:
    document = await owned_document(document_id, user, session)
    filename = safe_filename(payload.filename)
    if not filename.lower().endswith(".pdf"):
        filename += ".pdf"
    document.filename = filename
    await session.commit()
    await session.refresh(document)
    return document


@router.delete("/{document_id}", status_code=204)
async def delete_document(
    document_id: uuid.UUID, user: User = Depends(current_user), session: AsyncSession = Depends(get_session),
) -> Response:
    document = await owned_document(document_id, user, session)
    try:
        ObjectStorage().remove(document.object_key)
    except Exception:
        pass
    if document.original_object_key:
        try:
            ObjectStorage().remove(document.original_object_key)
        except Exception:
            pass
    await session.delete(document)
    await session.commit()
    return Response(status_code=204)


@router.get("/{document_id}/job", response_model=ProcessingJobResponse)
async def get_processing_job(
    document_id: uuid.UUID,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> ProcessingJob:
    await owned_document(document_id, user, session)
    job = await session.scalar(
        select(ProcessingJob).where(ProcessingJob.document_id == document_id).order_by(ProcessingJob.created_at.desc())
    )
    if job is None:
        raise HTTPException(status_code=404, detail="Processing job not found")
    return job


@router.get("/{document_id}/pages", response_model=list[DocumentPageResponse])
async def get_document_pages(
    document_id: uuid.UUID,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> list[DocumentPage]:
    await owned_document(document_id, user, session)
    result = await session.scalars(
        select(DocumentPage).where(DocumentPage.document_id == document_id).order_by(DocumentPage.page_number)
    )
    return list(result)


@router.get("/{document_id}/content")
async def download_document(
    document_id: uuid.UUID,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> StreamingResponse:
    document = await owned_document(document_id, user, session)
    data = ObjectStorage().download(document.object_key)
    disposition = f'inline; filename="{safe_filename(document.filename)}"'
    return StreamingResponse(BytesIO(data), media_type="application/pdf", headers={"Content-Disposition": disposition})


@router.get("/{document_id}/original")
async def download_original_document(
    document_id: uuid.UUID,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> StreamingResponse:
    document = await owned_document(document_id, user, session)
    key = document.original_object_key or document.object_key
    filename = document.original_filename or document.filename
    content_type = document.original_content_type or document.content_type
    data = ObjectStorage().download(key)
    return StreamingResponse(
        BytesIO(data),
        media_type=content_type,
        headers={"Content-Disposition": f'attachment; filename="{safe_filename(filename)}"'},
    )


@router.get("/{document_id}/thumbnail")
async def document_thumbnail(
    document_id: uuid.UUID,
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> Response:
    document = await owned_document(document_id, user, session)
    pdf = fitz.open(stream=ObjectStorage().download(document.object_key), filetype="pdf")
    try:
        if not pdf.page_count:
            raise HTTPException(status_code=422, detail="The PDF has no pages")
        image = pdf[0].get_pixmap(matrix=fitz.Matrix(.55, .55), alpha=False).tobytes("png")
    finally:
        pdf.close()
    return Response(
        content=image,
        media_type="image/png",
        headers={"Cache-Control": "private, max-age=3600"},
    )


@router.post("", response_model=DocumentResponse, status_code=status.HTTP_201_CREATED)
async def upload_document(
    file: UploadFile = File(...),
    workspace_id: str | None = Form(default=None),
    user: User = Depends(current_user),
    session: AsyncSession = Depends(get_session),
) -> Document:
    settings = get_settings()
    today = datetime.now(UTC).date()
    daily_user_uploads = await session.scalar(
        select(func.count(Document.id)).where(
            Document.owner_id == user.id,
            func.date(Document.created_at) == today,
        )
    )
    if (daily_user_uploads or 0) >= settings.daily_upload_limit_per_user:
        raise HTTPException(status_code=429, detail="Daily upload limit reached")
    daily_global_uploads = await session.scalar(
        select(func.count(Document.id)).where(func.date(Document.created_at) == today)
    )
    if (daily_global_uploads or 0) >= settings.global_daily_upload_limit:
        raise HTTPException(status_code=429, detail="The preview has reached its daily upload limit")
    document_count = await session.scalar(select(func.count(Document.id)).where(Document.owner_id == user.id))
    if (document_count or 0) >= settings.max_documents_per_user:
        raise HTTPException(status_code=422, detail=f"Document limit reached ({settings.max_documents_per_user})")
    original_filename = safe_filename(file.filename or "document.pdf")
    suffix = original_filename.lower().rsplit(".", 1)[-1] if "." in original_filename else ""
    generic_upload_types = {None, "", "application/octet-stream", "application/zip"}
    is_pdf = suffix == "pdf" and (file.content_type in {"application/pdf", "application/x-pdf", *generic_upload_types} or not file.content_type)
    is_image = suffix in {"png", "jpg", "jpeg", "webp"} and (file.content_type in {*IMAGE_CONTENT_TYPES, *generic_upload_types} or not file.content_type)
    is_docx = suffix == "docx" and (file.content_type in {DOCX_CONTENT_TYPE, *generic_upload_types} or not file.content_type)
    is_pptx = suffix == "pptx" and (file.content_type in {PPTX_CONTENT_TYPE, *generic_upload_types} or not file.content_type)
    is_text = suffix in {"txt", "md", "markdown", "rtf"} and (file.content_type in {*TEXT_CONTENT_TYPES, *generic_upload_types} or not file.content_type)
    if not is_pdf and not is_image and not is_docx and not is_pptx and not is_text:
        raise HTTPException(status_code=415, detail="Supported sources are PDF, DOCX, PPTX, Markdown, text, RTF, PNG, JPEG, and WebP")
    data = await file.read(settings.max_file_size_mb * 1024 * 1024 + 1)
    if len(data) > settings.max_file_size_mb * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"File exceeds {settings.max_file_size_mb} MB")
    if is_pdf and not data.startswith(b"%PDF-"):
        raise HTTPException(status_code=422, detail="The file is not a valid PDF")
    from app.deliverables import activity, ensure_personal_workspace, workspace_access

    ws_uuid: uuid.UUID | None = None
    if workspace_id and str(workspace_id).strip() and str(workspace_id).strip().lower() not in {"null", "undefined"}:
        try:
            ws_uuid = uuid.UUID(str(workspace_id).strip())
        except (ValueError, AttributeError):
            ws_uuid = None

    if ws_uuid is not None:
        workspace, _ = await workspace_access(ws_uuid, user, session, {"owner", "editor"})
    else:
        workspace = await ensure_personal_workspace(user, session)

    source_sha256 = hashlib.sha256(data).hexdigest()
    duplicate = await session.scalar(select(Document).where(
        Document.owner_id == user.id, Document.workspace_id == workspace.id, Document.source_sha256 == source_sha256
    ))
    if duplicate is not None:
        raise HTTPException(status_code=409, detail={"message": "This source is already in your workspace", "document_id": str(duplicate.id)})
    display_title = None
    filename = original_filename
    original_data = data
    if is_image:
        try:
            data = await run_in_threadpool(image_to_pdf, data)
        except ValueError as exc:
            raise HTTPException(status_code=422, detail=str(exc)) from exc
        filename = f"{original_filename.rsplit('.', 1)[0]}.pdf"
        display_title = original_filename
    elif not is_pdf:
        try:
            data = await run_in_threadpool(source_to_pdf, data, suffix, original_filename)
        except ValueError as exc:
            raise HTTPException(status_code=422, detail=str(exc)) from exc
        filename = f"{original_filename.rsplit('.', 1)[0]}.pdf"
        display_title = original_filename
    document_id = uuid.uuid4()
    object_key = f"{user.id}/{document_id}/{filename}"
    original_object_key = None if is_pdf else f"{user.id}/{document_id}/original/{original_filename}"
    storage = ObjectStorage()
    try:
        storage.upload_pdf(object_key, data)
        if original_object_key:
            storage.upload(original_object_key, original_data, file.content_type or "application/octet-stream")
    except Exception as exc:
        try:
            storage.remove(object_key)
        except Exception:
            pass
        if original_object_key:
            try:
                storage.remove(original_object_key)
            except Exception:
                pass
        raise HTTPException(status_code=503, detail=f"Document storage error: {str(exc)}") from exc
    document = Document(
        id=document_id,
        owner_id=user.id,
        workspace_id=workspace.id,
        filename=filename,
        object_key=object_key,
        content_type="application/pdf",
        size_bytes=len(data),
        display_title=display_title,
        original_filename=original_filename if original_object_key else None,
        original_object_key=original_object_key,
        original_content_type=file.content_type if original_object_key else None,
        source_sha256=source_sha256,
    )
    session.add(document)
    await activity(session, workspace.id, user.id, "source.uploaded", "document", document.id, {"title": display_title or filename})
    job = ProcessingJob(
        document_id=document_id,
        owner_id=user.id,
        operation="document_processing",
        status=JobStatus.QUEUED,
        progress=0,
    )
    session.add(job)
    await session.commit()
    await session.refresh(document)
    from app.tasks import process_document

    try:
        task = process_document.delay(str(document.id))
        job.task_id = task.id
        await session.commit()
    except Exception as exc:
        document.status = DocumentStatus.FAILED
        document.error_message = f"Processing worker offline: {str(exc)}"
        job.status = JobStatus.FAILED
        job.error_message = document.error_message
        await session.commit()
    return document


@router.post("/{document_id}/retry", response_model=ProcessingJobResponse)
async def retry_document(
    document_id: uuid.UUID, user: User = Depends(current_user), session: AsyncSession = Depends(get_session),
) -> ProcessingJob:
    document = await owned_document(document_id, user, session)
    if document.status != DocumentStatus.FAILED:
        raise HTTPException(status_code=409, detail="Only failed documents can be retried")
    job = ProcessingJob(
        document_id=document.id,
        owner_id=user.id,
        operation="document_processing",
        status=JobStatus.QUEUED,
        progress=0,
    )
    document.status = DocumentStatus.UPLOADED
    document.error_message = None
    session.add(job)
    from app.deliverables import activity, ensure_personal_workspace
    workspace = await ensure_personal_workspace(user, session)
    await activity(session, workspace.id, user.id, "source.processing_retried", "document", document.id, {"title": document.display_title or document.filename})
    await session.commit()
    await session.refresh(job)
    from app.tasks import process_document
    try:
        task = process_document.delay(str(document.id))
        job.task_id = task.id
        await session.commit()
    except Exception as exc:
        job.status = JobStatus.FAILED
        job.error_message = f"Processing queue unavailable: {str(exc)}"
        document.status = DocumentStatus.FAILED
        document.error_message = job.error_message
        await session.commit()
        raise HTTPException(status_code=503, detail=job.error_message) from exc
    return job
